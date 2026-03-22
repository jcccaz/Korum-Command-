# CONFIDENTIAL - TRADE SECRET
# Proprietary KorumOS source code. Access is limited to authorized personnel
# and collaborators operating under written confidentiality obligations.
#
# -------------------------------------------------------------------
# KORUM V2: SEQUENTIAL COUNCIL ENGINE (Python Implementation)
# -------------------------------------------------------------------

import os
import re
import json
from datetime import datetime

from llm_core import call_openai_gpt4

# Note: We rely on the core functions. If they are not available as discrete imports from app.py
# (due to circular dependency), we now import them from llm_core.
# I will assume `llm_core.py` was created successfully in the previous step.
# Import the other providers:
from llm_core import call_anthropic_claude, call_google_gemini, call_perplexity, call_local_llm, call_mistral_api

# --- WORKFLOW DNA REGISTRY (The Elite Logic Layer) ---
WORKFLOW_DNA = {
    "WAR_ROOM": {
        "goal": "Immediate tactical action and crisis containment.",
        "tone": "Aggressive, direct, and zero-fluff.",
        "risk_bias": "Conservative (Minimize immediate damage)",
        "time_horizon": "0–72 hours",
        "posture": "Tactical Commander",
        "output_structure": ["Situation", "Threat", "Immediate Action", "Resource Allocation", "Escalation Path"]
    },
    "RESEARCH": {
        "goal": "Clear decision with supporting analysis and actionable next steps.",
        "tone": "Direct, structured, and decision-ready.",
        "risk_bias": "Balanced",
        "time_horizon": "Long-term strategic",
        "posture": "Decision Commander",
        "output_structure": ["Executive Summary", "Key Signals", "System Context", "Scenario Analysis", "Critical Challenges", "Tradeoffs", "Decision", "Action Priorities", "Confidence"]
    },
    "FINANCE": {
        "goal": "Economic viability and downside protection.",
        "tone": "Analytical, precise, and detached.",
        "risk_bias": "Downside-aware",
        "time_horizon": "Scenario-based",
        "posture": "CFO / Auditor",
        "output_structure": ["Cost", "Revenue Impact", "Sensitivity Table", "Worst-case Scenario", "ROI Summary"]
    },
    "LEGAL": {
        "goal": "Exposure reduction and regulatory compliance.",
        "tone": "Formal, rigorous, and protective.",
        "risk_bias": "Zero-risk / Protective",
        "time_horizon": "Indefinite",
        "posture": "General Counsel",
        "output_structure": ["Regulatory Exposure", "Contractual Impact", "Risk Mitigation", "Recommended Posture"]
    },
    "QUANTUM_SECURITY": {
        "goal": "Assess cryptographic vulnerabilities, enforce Zero Trust architecture, and map to strict government compliance (NIST 800-207, FedRAMP, CMMC).",
        "tone": "Authoritative, highly technical, and uncompromising on security.",
        "risk_bias": "Zero-trust / Absolute Security",
        "time_horizon": "Long-term (Post-Quantum Readiness)",
        "posture": "Chief Information Security Officer (CISO) & Cryptographer",
        "output_structure": ["Threat Landscape", "Cryptographic Vulnerabilities", "Zero Trust Controls", "Compliance Mapping (NIST/FedRAMP)", "Mitigation Architecture"]
    },
    "MEDICAL": {
        "goal": "Evidence-based clinical assessment with patient safety as the absolute priority.",
        "tone": "Precise, cautious, and citation-driven.",
        "risk_bias": "Conservative (Do No Harm)",
        "time_horizon": "Immediate clinical + long-term outcomes",
        "posture": "Chief Medical Officer & Clinical Review Board",
        "output_structure": ["Clinical Assessment", "Evidence Review", "Differential Diagnosis", "Treatment Options", "Risk-Benefit Analysis"]
    },
    "CYBER": {
        "goal": "Identify active threats, map attack surfaces, and provide actionable defense recommendations.",
        "tone": "Direct, technical, and adversary-aware.",
        "risk_bias": "Assume breach / Adversarial",
        "time_horizon": "0-72 hours tactical + ongoing hardening",
        "posture": "Cyber Command / Red Team Lead",
        "output_structure": ["Threat Intelligence", "Attack Surface Analysis", "Active Indicators (IOCs)", "Mitigation Playbook", "Detection Rules"]
    },
    "DEFENSE": {
        "goal": "Strategic defense analysis with geopolitical awareness and operational security.",
        "tone": "Classified-brief style, precise, and mission-focused.",
        "risk_bias": "Threat-forward / Worst-case planning",
        "time_horizon": "Multi-domain (tactical to strategic)",
        "posture": "Defense Intelligence Analyst",
        "output_structure": ["Situation Assessment", "Force Disposition", "Threat Analysis", "Courses of Action", "Intelligence Gaps"]
    },
    "STARTUP": {
        "goal": "Validate business viability, identify product-market fit, and assess go-to-market readiness.",
        "tone": "Pragmatic, founder-direct, and metrics-driven.",
        "risk_bias": "Calculated risk / Move fast with data",
        "time_horizon": "0-18 months runway",
        "posture": "Startup Advisor / VC Partner",
        "output_structure": ["Market Opportunity", "Unit Economics", "Competitive Landscape", "Go-to-Market Strategy", "Funding & Runway"]
    },
    "AUDIT": {
        "goal": "Identify control failures, compliance gaps, and evidence discrepancies against applicable frameworks.",
        "tone": "Formal, objective, and finding-driven.",
        "risk_bias": "Skeptical / Trust but verify",
        "time_horizon": "Point-in-time assessment + remediation timeline",
        "posture": "Lead Auditor (SOC2 / ISO 27001 / NIST)",
        "output_structure": ["Scope & Methodology", "Control Findings", "Compliance Gap Analysis", "Evidence Assessment", "Remediation Priorities"]
    },
    "CREATIVE": {
        "goal": "Generate innovative concepts, narratives, and strategic messaging — then draft the actual deliverable.",
        "tone": "Bold, imaginative, and audience-aware.",
        "risk_bias": "Risk-tolerant / Push boundaries",
        "time_horizon": "Campaign-based",
        "posture": "Creative Director & Brand Strategist",
        "output_structure": ["Creative Concept", "Target Audience", "Messaging Framework", "Draft Deliverable", "Execution Plan"]
    },
    "SCIENCE": {
        "goal": "Evaluate hypotheses with methodological rigor and reproducibility standards.",
        "tone": "Academic, precise, and peer-review caliber.",
        "risk_bias": "Evidence-weighted / Null hypothesis default",
        "time_horizon": "Study-dependent",
        "posture": "Principal Investigator & Peer Reviewer",
        "output_structure": ["Hypothesis Evaluation", "Methodology Assessment", "Data Analysis", "Literature Context", "Conclusions & Limitations"]
    },
    "TECH": {
        "goal": "Evaluate technical architecture, scalability, and engineering trade-offs.",
        "tone": "Engineer-direct, pragmatic, and systems-aware.",
        "risk_bias": "Reliability-first",
        "time_horizon": "Build cycle (sprint to quarterly)",
        "posture": "CTO / Principal Engineer",
        "output_structure": ["Architecture Assessment", "Technical Trade-offs", "Scalability Analysis", "Implementation Plan", "Technical Debt & Risks"]
    },
    "INTEL": {
        "goal": "Produce finished intelligence products with source evaluation and confidence grading.",
        "tone": "Intelligence community standard, structured, and source-attributed.",
        "risk_bias": "Analytic confidence calibration",
        "time_horizon": "Current intelligence + strategic forecast",
        "posture": "Senior Intelligence Analyst",
        "output_structure": ["Key Judgments", "Source Evaluation", "Analytic Confidence", "Alternative Hypotheses", "Intelligence Gaps"]
    },
    "SYSTEM": {
        "goal": "Diagnose system behavior, optimize performance, and ensure operational integrity.",
        "tone": "Diagnostic, systematic, and root-cause focused.",
        "risk_bias": "Stability-first",
        "time_horizon": "Immediate resolution + preventive measures",
        "posture": "Systems Engineer / SRE Lead",
        "output_structure": ["System Status", "Root Cause Analysis", "Performance Metrics", "Resolution Steps", "Monitoring Recommendations"]
    },
    # --- ALIASES TO BRIDGE UI NAMES ---
    "CREATIVE_COUNCIL": { "alias": "CREATIVE" },
    "STARTUP_LAUNCH": { "alias": "STARTUP" },
    "CODE_AUDIT": { "alias": "AUDIT" },
    "CYBER_COMMAND": { "alias": "CYBER" },
    "DEFENSE_COUNCIL": { "alias": "DEFENSE" },
    "INTEL_BRIEF": { "alias": "INTEL" },
    "SCIENCE_PANEL": { "alias": "SCIENCE" },
    "SOCIAL_POST": {
        "goal": "Draft high-impact technical narratives that positions the author as a builder-authority and stops the scroll.",
        "tone": "Authentic, technical but accessible, and zero-fluff.",
        "risk_bias": "Trust-focused (Build in public)",
        "time_horizon": "Immediate impact",
        "posture": "Founder-Builder & Storyteller",
        "output_structure": ["Target Audience", "Core Hook", "Main Draft", "Alternate Versions", "Engagement Logic"]
    },
    "EOM_STATEMENT": {
        "goal": "Generate a professional, structured End-of-Month (EOM) financial report with zero fluff.",
        "tone": "Formal, precise, and accounting-standard aligned.",
        "risk_bias": "Conservative (Downside-aware)",
        "time_horizon": "Monthly / Trailing 12-month",
        "posture": "Chief Financial Officer (CFO)",
        "output_structure": ["P&L Statement", "Balance Sheet Snapshot", "Cash Flow & Burn Rate", "Runway & Forecast", "Financial Priorities"]
    },
    "PORTFOLIO_BUILDER": {
        "goal": "Build a complete, actionable investment portfolio with specific tickers, allocations, entry ranges, and price targets. No disclaimers. No hedging. Picks with numbers.",
        "tone": "Aggressive, data-driven, hedge fund floor. Direct calls only.",
        "risk_bias": "Risk-calibrated — position size to the conviction level, not to comfort.",
        "time_horizon": "12-month primary / 3-year thesis",
        "posture": "Hedge Fund Portfolio Manager & Investment Committee",
        "output_structure": ["Macro Setup", "Screened Candidates", "Red Team Challenge", "Portfolio Architecture", "Final Portfolio — Ready to Execute"]
    }
}

# --- FINAL ARBITER: METRIC VOCABULARY (Section 8a Enforcement) ---
# Maps common financial labels to canonical metric names.
# The arbiter uses this to recognise the same concept across providers.
METRIC_ALIASES = {
    "revenue":    ["revenue", "total revenue", "gross revenue", "top-line", "top line", "sales", "total sales", "net sales"],
    "costs":      ["costs", "total costs", "expenses", "total expenses", "operating costs", "opex", "cogs",
                   "cost of goods sold", "operating expenses", "total operating costs"],
    "profit":     ["profit", "net profit", "net income", "earnings", "operating profit", "ebitda",
                   "gross profit", "operating income", "net earnings", "bottom line", "bottom-line"],
    "margin":     ["margin", "profit margin", "net margin", "gross margin", "operating margin"],
    "roi":        ["roi", "return on investment"],
    "growth":     ["growth", "growth rate", "yoy growth", "year-over-year growth", "cagr", "revenue growth"],
    "market_cap": ["market cap", "market capitalization", "valuation"],
    "debt":       ["debt", "total debt", "liabilities", "total liabilities"],
    "assets":     ["assets", "total assets"],
    "equity":     ["equity", "total equity", "shareholders equity", "net worth", "book value"],
    "cash_flow":  ["cash flow", "free cash flow", "fcf", "operating cash flow"],
    "burn_rate":  ["burn rate", "monthly burn", "cash burn", "net burn"],
    "runway":     ["runway", "cash runway", "months of runway"],
}

# Arithmetic identities the arbiter validates.
# Format: (result_metric, operator, operand_a, operand_b)
MATH_IDENTITIES = [
    ("profit",  "subtract",    "revenue", "costs"),      # profit = revenue - costs
    ("margin",  "divide_pct",  "profit",  "revenue"),    # margin = profit / revenue * 100
    ("roi",     "divide_pct",  "profit",  "costs"),      # roi    = profit / costs   * 100
    ("equity",  "subtract",    "assets",  "debt"),       # equity = assets - debt
]

# --- DETERMINISTIC PIPELINE REGISTRY (Locked Production Paths) ---
# Maps workflows to a fixed list of [provider-role] steps.
# This ensures that for any given workflow, the same AI roles and providers
# are triggered in the same sequence every time.
# Format: [Phase 0, Phase 1, Phase 2, Phase 3, Phase 4]
WORKFLOW_STEPS = {
    "WAR_ROOM": ["openai-commander", "anthropic-tactician", "google-analyst", "perplexity-scout", "mistral-validator"],
    "RESEARCH": ["openai-analyst", "anthropic-architect", "google-critic", "perplexity-scout", "mistral-validator"],
    "FINANCE": ["openai-cfo", "anthropic-auditor", "google-quant", "perplexity-researcher", "mistral-compliance"],
    "LEGAL": ["openai-counsel", "anthropic-analyst", "google-critic", "perplexity-researcher", "mistral-validator"],
    "QUANTUM_SECURITY": ["openai-ciso", "anthropic-cryptographer", "google-security_auditor", "perplexity-threat_intel", "mistral-compliance"],
    "MEDICAL": ["openai-physician", "anthropic-specialist", "google-researcher", "perplexity-clinical_scout", "mistral-compliance"],
    "CYBER": ["openai-threat_analyst", "anthropic-architect", "google-red_team", "perplexity-threat_intel", "mistral-incident_lead"],
    "DEFENSE": ["openai-intelligence_officer", "anthropic-strategist", "google-adversarial_analyst", "perplexity-field_scout", "mistral-commander"],
    "STARTUP": ["openai-founder", "anthropic-advisor", "google-critic", "perplexity-market_scout", "mistral-strategist"],
    "AUDIT": ["openai-auditor", "anthropic-compliance_lead", "google-skeptic", "perplexity-fact_checker", "mistral-validator"],
    "CREATIVE": ["openai-creative_director", "anthropic-brand_strategist", "google-critic", "perplexity-cultural_scout", "mistral-copywriter"],
    "SCIENCE": ["openai-principal_investigator", "anthropic-methodologist", "google-peer_reviewer", "perplexity-literature_scout", "mistral-statistician"],
    "TECH": ["openai-cto", "anthropic-architect", "google-failure_analyst", "perplexity-tech_scout", "mistral-validator"],
    "INTEL": ["openai-collection_manager", "anthropic-analyst", "google-counterintel", "perplexity-osint_scout", "mistral-intelligence_officer"],
    "SYSTEM": ["openai-sre_lead", "anthropic-architect", "google-diagnostic_analyst", "perplexity-metrics_scout", "mistral-validator"],
    "SOCIAL_POST": ["openai-narrative_strategist", "anthropic-voice_analyst", "google-editor", "perplexity-trend_scout", "mistral-content_lead"],
    "EOM_STATEMENT": ["openai-forensic_accountant", "anthropic-financial_analyst", "google-cfo", "perplexity-variance_analyst", "mistral-auditor"],
    "PORTFOLIO_BUILDER": ["openai-macro_strategist", "anthropic-quant_analyst", "google-short_seller", "perplexity-risk_manager", "mistral-cio"],
}



# --- STRATEGIC SIGNAL TAG REGISTRY ---
# Use these tags to force immediate scan-pattern recognition.
# [CRITICAL] - Red Reserve (Use sparingly, <10% of total tags)
# [ACTION REQUIRED] - Amber Default
# [VERIFIED] - Amber Default
# [RISK] - Amber Default
# [WARNING] - Amber Default
# [REJECTED] - High-visibility callout of failed truth-checks.

# --- PHASE DIRECTIVES (Global Registry) ---
# Directives are domain-adaptive: the structure stays constant but the
# language flexes to match the actual query topic.
PHASE_DIRECTIVES = {
    0: {
        "title": "ANALYST — Problem Framing",
        "instruction": (
            "You are the ANALYST. Your job is to frame the problem and establish the factual baseline.\n\n"
            "MUST INCLUDE:\n"
            "- Core entities and their roles in the situation\n"
            "- Observed conditions (not 'Verified Facts' — use what is visible, allow light inference)\n"
            "- 2-4 cause-oriented hypotheses for the situation\n"
            "- Key risks and their directional severity\n"
            "- Unknowns and data gaps (but do NOT overuse 'data not provided' — infer where reasonable)\n\n"
            "MUST NOT INCLUDE:\n"
            "- Scenarios, recommendations, or decisions\n"
            "- Action plans or implementation steps\n"
            "- Confidence scores or final judgments\n"
            "- Tool/vendor recommendations\n\n"
            "TONE: Neutral, structured, fact-forward. You are mapping the terrain, not choosing a direction."
        )
    },
    1: {
        "title": "ARCHITECT — Scenario Modeling",
        "instruction": (
            "You are the ARCHITECT. Your job is to model strategic scenarios from the Analyst's baseline.\n\n"
            "MUST INCLUDE:\n"
            "- Most Plausible scenario with progression timeline\n"
            "- Most Dangerous scenario with escalation triggers\n"
            "- Strategic Opportunity scenario (upside lens, not just downside)\n"
            "- Implications for each scenario (one-sentence impact statement)\n\n"
            "MUST NOT INCLUDE:\n"
            "- Extreme stacked worst-case outcomes in a single scenario\n"
            "- Unsupported precise numbers (e.g., '37.2% revenue loss') — use directional ranges\n"
            "- Blanket 'high confidence' claims without grounding\n"
            "- Final decisions or action plans\n"
            "- Tool/vendor lists or implementation details\n\n"
            "TONE: Strategic, not alarmist. Keep the opportunity lens. "
            "Avoid synthetic precision like '25-40% churn' unless the input data supports it."
        )
    },
    2: {
        "title": "CRITIC — Assumption Challenge",
        "instruction": (
            "You are the CRITIC. Your job is to stress-test the Architect's scenarios and assumptions.\n\n"
            "MUST INCLUDE:\n"
            "- Prioritized challenges (High Impact first, then Secondary)\n"
            "- WHY each challenge matters (not just what is wrong)\n"
            "- What specific data or validation would resolve each challenge\n"
            "- A synthesis line: do uncertainties outweigh the risk of inaction?\n\n"
            "MUST NOT INCLUDE:\n"
            "- Final decisions (you challenge, you do NOT decide)\n"
            "- 'DECISION: No' or 'DECISION: Halt' labels — use 'Challenge' / 'Constraint' / 'Validation Required'\n"
            "- Analysis of the report itself (no references to 'the Executive Summary section' or 'Key Signals')\n"
            "- Recommendations, action plans, or implementation steps\n"
            "- New scenarios or strategic options\n\n"
            "TONE: Constructive adversarial. You are the quality gate, not the decision-maker."
        )
    },
    3: {
        "title": "INTEGRATOR — Decision Synthesis",
        "instruction": (
            "You are the INTEGRATOR. You are the SOLE decision authority. "
            "Resolve tensions between Analyst, Architect, and Critic into ONE clear decision.\n\n"
            "MUST INCLUDE:\n"
            "- The final decision (one clear sentence)\n"
            "- Rationale (3-5 bullet points)\n"
            "- Key tradeoffs considered\n"
            "- High-level action priorities (2-3 per time horizon: Immediate, Near-Term, Mid-Term)\n"
            "- Confidence assessment with reasoning\n\n"
            "MUST NOT INCLUDE:\n"
            "- Step-by-step procedures or project-management timelines\n"
            "- Methodology tutorials or framework explanations\n"
            "- Micro-level execution plans (no 'Day 1-3' schedules)\n"
            "- Tool/vendor lists unless absolutely essential to the decision\n"
            "- Detailed staffing, budgets, or operational checklists\n\n"
            "TONE: Decisive executive. You own the call. This is NOT an ops playbook — "
            "it is a decision with enough context to act on."
        )
    },
    4: {
        "title": "COMPOSER — Executive Narrative",
        "instruction": (
            "You are the REPORT COMPOSER. Your job is to unify ALL prior analysis into one polished executive artifact.\n\n"
            "MUST:\n"
            "- Write in ONE unified senior consultant voice\n"
            "- Lead with the decision in the Executive Summary\n"
            "- Follow the output_structure exactly (all sections, in order)\n"
            "- Keep actions high-level (max 2-3 per time horizon)\n"
            "- Use directional confidence aligned with actual data support\n\n"
            "MUST NOT:\n"
            "- Expose node numbers, model names, or provider identities\n"
            "- Use internal tags like [VERIFIED], [CRITICAL], [ACTION REQUIRED]\n"
            "- Include citations ([1][2][3]) in body text\n"
            "- Create day-by-day timelines or tool/vendor spam\n"
            "- Write in chat transcript or engineering SOP style\n"
            "- Say 'this model said' or 'the council found' — say 'analysis shows'\n\n"
            "VOICE: Senior strategy consultant delivering a decision memo to leadership."
        )
    }
}

FINANCE_PHASE_DIRECTIVES = {
    0: {
        "title": "INTAKE — Financial Extraction",
        "instruction": "Extract every financial figure from the source material. Build structured tables for Revenue, CAPEX, and OPEX."
    },
    1: {
        "title": "ANALYSIS — Financial Modeling",
        "instruction": "Model unit economics, burn rate, and sensitivity (±10%, ±25%). Use calculations, not opinions."
    },
    2: {
        "title": "STRESS TEST — Assumption Challenge",
        "instruction": "Skeptically challenge financial assumptions. What if costs are 2x? What if revenue is 0.5x?"
    },
    3: {
        "title": "COMPLIANCE & FRAMEWORKS — Regulatory Mapping",
        "instruction": "Map to accounting standards (GAAP/IFRS) and tax obligations."
    },
    4: {
        "title": "VERDICT — Full Financial Decision Package",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the CFO / LEAD AUDITOR. Deliver the complete financial decision package — NOT a summary paragraph. "
            "Structure as:\n"
            "  1) **DECISION** — Go / No-Go / Conditional (one line, bold)\n"
            "  2) **ROI TABLE** — Markdown table: [Scenario, ROI %, Payback Period, Key Assumption]\n"
            "  3) **SENSITIVITY MATRIX** — Base case / Bull case / Bear case with % outcomes\n"
            "  4) **TOP 3 RISKS** — Each with a specific mitigation action\n"
            "  5) **KEY SUCCESS METRICS** — The exact numbers that define success in 90 days\n"
            "STRICT RULE: Every number must come from prior phase analysis. No prose opinions. Tables and lists only."
        )
    }
}



WAR_ROOM_PHASE_DIRECTIVES = {
    0: { "title": "SITUATION REPORT", "instruction": "Establish ground truth in 60 seconds: What, Who, When, Unknowns." },
    1: { "title": "THREAT ASSESSMENT", "instruction": "Map blast radius and cascade risks for next 72 hours." },
    2: { "title": "IMMEDIATE ACTION", "instruction": "Containment only. STOP actions and PROTECT actions." },
    3: { "title": "RESOURCE ALLOCATION", "instruction": "Assign owners, tools, and budget to the containment plan." },
    4: { "title": "COMMANDER'S DECISION BRIEF", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the INCIDENT COMMANDER. Output the complete decision brief — ready to brief upward in 60 seconds. "
        "Structure as:\n"
        "  1) **SITREP** — Current status in 2 sentences\n"
        "  2) **DECISION REQUIRED** — The single yes/no choice that must be made NOW\n"
        "  3) **RECOMMENDED COA** — Your chosen course of action with rationale\n"
        "  4) **ESCALATION TRIPWIRES** — Markdown table: [Trigger, Threshold, Who Owns It, Action]\n"
        "  5) **DE-ESCALATION CRITERIA** — What does 'resolved' look like? Exact metrics.\n"
        "STRICT RULE: No new analysis. Compile and command. This brief must be copy-pasteable to a Slack channel or exec email."
    ) }
}

LEGAL_PHASE_DIRECTIVES = {
    0: { "title": "INTAKE — Exposure Mapping", "instruction": "Inventory parties, instruments, and legal trigger events." },
    1: { "title": "RISK ANALYSIS", "instruction": "Score every exposure with statute/case citations and jurisdiction." },
    2: { "title": "ADVERSARIAL REVIEW", "instruction": "Attack the position as opposing counsel. Find the fatal flaw." },
    3: { "title": "MITIGATION", "instruction": "Provide specific contract language and structural changes to reduce risk." },
    4: { "title": "LEGAL POSTURE — Full Opinion", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the GENERAL COUNSEL. Deliver the complete legal opinion package. Structure as:\n"
        "  1) **VERDICT** — Proceed / Do Not Proceed / Proceed With Conditions (bold, one line)\n"
        "  2) **RISK REGISTER** — Markdown table: [Risk, Severity (H/M/L), Jurisdiction, Statute/Case, Mitigation]\n"
        "  3) **NON-NEGOTIABLE PREREQUISITES** — Numbered list of conditions that MUST be satisfied before proceeding\n"
        "  4) **RECOMMENDED CONTRACT LANGUAGE** — Specific clause text or redline, ready to insert\n"
        "  5) **MONITORING TRIGGERS** — Events that require immediate legal review\n"
        "STRICT RULE: Cite actual statutes, regulations, or case names. No generic legal hedging."
    ) }
}

MEDICAL_PHASE_DIRECTIVES = {
    0: { "title": "INTAKE — Clinical Fact Extraction", "instruction": "Extract symptoms, history, and labs. Identify missing data." },
    1: { "title": "DIFFERENTIAL", "instruction": "Rank candidate diagnoses with ICD-10 codes and evidence grades." },
    2: { "title": "CHALLENGE — Blind Spots", "instruction": "Find rare or dangerous conditions missed. Audit for cognitive bias." },
    3: { "title": "TREATMENT PLAN", "instruction": "Evidence-based interventions with dosage ranges and contraindications." },
    4: { "title": "CLINICAL DECISION BRIEF", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the CHIEF MEDICAL OFFICER. Deliver the complete clinical decision brief. Structure as:\n"
        "  1) **PRIMARY DIAGNOSIS** — ICD-10 code, confidence level, and one-sentence rationale\n"
        "  2) **TREATMENT PLAN TABLE** — Markdown table: [Intervention, Dosage/Protocol, Evidence Grade, Contraindications]\n"
        "  3) **DIFFERENTIAL WATCHLIST** — Top 2 ruled-out diagnoses and what would change the verdict\n"
        "  4) **DRUG INTERACTION FLAGS** — Any interactions that require immediate attention\n"
        "  5) **INFORMED CONSENT CHECKLIST** — Bullet list of what the patient must be told before proceeding\n"
        "STRICT RULE: Cite clinical guidelines (e.g. AHA, WHO, UpToDate). No general health advice."
    ) }
}

QUANTUM_SECURITY_PHASE_DIRECTIVES = {
    0: { "title": "INVENTORY", "instruction": "Map algorithms, PKI, and quantum exposure (Shor's/Grover's impact)." },
    1: { "title": "VULNERABILITY", "instruction": "Score vs classical (BEAST/POODLE) and emerging quantum threats." },
    2: { "title": "ZERO TRUST AUDIT", "instruction": "Audit vs NIST 800-207 pillars (Identity/Device/Network/App/Data)." },
    3: { "title": "COMPLIANCE MAPPING", "instruction": "Map current gaps to FedRAMP, CMMC, and NIST IR 8413." },
    4: { "title": "PQC MIGRATION ROADMAP", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the CISO. Deliver the complete Post-Quantum Cryptography migration plan. Structure as:\n"
        "  1) **MIGRATION PRIORITY TABLE** — Markdown table: [System/Algorithm, Quantum Risk, Replace With, Timeline, Owner]\n"
        "  2) **PHASE 1 (0-90 days)** — Immediate actions: inventory, disable deprecated algorithms\n"
        "  3) **PHASE 2 (90-180 days)** — Migration to ML-KEM (key encapsulation) and ML-DSA (signatures)\n"
        "  4) **PHASE 3 (180+ days)** — Full PQC posture, audit, and compliance sign-off\n"
        "  5) **COMPLIANCE CHECKLIST** — NIST IR 8413, FedRAMP, CMMC specific sign-off items\n"
        "STRICT RULE: Name specific algorithms (e.g. Kyber-1024, Dilithium-3). No generic 'upgrade encryption' language."
    ) }
}

DEFENSE_PHASE_DIRECTIVES = {
    0: { "title": "OPERATIONAL PICTURE", "instruction": "Force disposition: Blue Force, Red Force, Grey Force, Key Terrain." },
    1: { "title": "THREAT ASSESSMENT", "instruction": "Model MLCOA (Most Likely) and MDCOA (Most Dangerous) adversary moves." },
    2: { "title": "RED CELL", "instruction": "Attack our plan from the adversary perspective. Find the feint." },
    3: { "title": "COURSES OF ACTION", "instruction": "Friendly options: Scheme of maneuver, requirements, and decision points." },
    4: { "title": "COMMANDER'S ESTIMATE", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the COMMANDING OFFICER. Deliver the complete commander's estimate. Structure as:\n"
        "  1) **MISSION** — Restated in one sentence\n"
        "  2) **SELECTED COA** — The chosen course of action (bold) with a 2-sentence rationale\n"
        "  3) **EXECUTION TABLE** — Markdown table: [Phase, Task, Unit/Owner, Timeline, Success Criteria]\n"
        "  4) **ACCEPTED RISKS** — What risks are being accepted and why\n"
        "  5) **DECISION POINTS** — Specific triggers that would force a COA change\n"
        "STRICT RULE: One recommendation only. No hedging. State risk acceptance explicitly."
    ) }
}

CYBER_PHASE_DIRECTIVES = {
    0: { "title": "TRIAGE — IOCs", "instruction": "Extract indicators (IPs, hashes) and build the incident timeline." },
    1: { "title": "ATTACK ANALYSIS", "instruction": "Map to MITRE ATT&CK techniques and kill chain stages." },
    2: { "title": "ADVERSARIAL INTENT", "instruction": "Predict objectives: exfiltration, ransomware, or persistence?" },
    3: { "title": "CONTAINMENT", "instruction": "Immediate stop/remove actions and evidence preservation." },
    4: { "title": "HARDENING PLAYBOOK", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the CYBER DEFENSE LEAD. Deliver the complete post-incident hardening playbook. Structure as:\n"
        "  1) **IMMEDIATE ACTIONS (24h)** — Numbered list: specific firewall rules, account disables, or patch installs\n"
        "  2) **SIEM RULES** — Exact detection logic (e.g. sigma rule pseudocode or query) for each identified TTPs\n"
        "  3) **EDR POLICY CHANGES** — Specific policy updates with product names (CrowdStrike, Defender, etc.)\n"
        "  4) **NIST CONTROL MAP** — Markdown table: [MITRE Technique, NIST 800-53 Control, Status, Owner]\n"
        "  5) **30-DAY VALIDATION PLAN** — How you confirm the hardening worked\n"
        "STRICT RULE: Be tool-specific. No generic 'patch your systems' advice."
    ) }
}

INTEL_PHASE_DIRECTIVES = {
    0: { "title": "SOURCE INVENTORY", "instruction": "Rate sources by reliability (NATO STANAG 2022) and info certainty." },
    1: { "title": "ANALYTIC LINE", "instruction": "Numbered judgments with ICD 203 confidence levels." },
    2: { "title": "ALT ANALYSIS", "instruction": "Structured analytic techniques: Analysis of Competing Hypotheses (ACH)." },
    3: { "title": "IMPLICATIONS", "instruction": "So what? Policy implications and opportunity windows." },
    4: { "title": "FINISHED INTELLIGENCE PRODUCT", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the SENIOR INTELLIGENCE OFFICER. Deliver the complete finished intelligence product. Structure as:\n"
        "  1) **KEY JUDGMENTS** — Numbered list, each starting with an ICD 203 confidence label (High/Moderate/Low)\n"
        "  2) **BLUF** — Bottom Line Up Front in 2 sentences max\n"
        "  3) **EVIDENCE SUMMARY TABLE** — Markdown table: [Source, Reliability, Info Certainty, Key Finding]\n"
        "  4) **ALTERNATIVE HYPOTHESIS** — The most credible competing explanation and why it was downgraded\n"
        "  5) **COLLECTION GAPS** — Specific intelligence gaps and recommended collection platforms/methods\n"
        "STRICT RULE: Every judgment must have a confidence level. No ungraded assertions."
    ) }
}

SCIENCE_PHASE_DIRECTIVES = {
    0: { "title": "HYPOTHESIS FRAMING", "instruction": "Testable research question and null hypothesis." },
    1: { "title": "EVIDENCE ANALYSIS", "instruction": "Evaluate study design (RCT/cohort) and statistical power." },
    2: { "title": "PEER REVIEW", "instruction": "Methodology challenge: confounding variables and selection bias." },
    3: { "title": "SYNTHESIS", "instruction": "What the evidence says: Verdict on hypothesis vs practical significance." },
    4: { "title": "RESEARCH CONCLUSION & AGENDA", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the PRINCIPAL INVESTIGATOR. Deliver the complete research conclusion package. Structure as:\n"
        "  1) **VERDICT ON HYPOTHESIS** — Supported / Refuted / Inconclusive with confidence grade (p-value or effect size if available)\n"
        "  2) **SUMMARY OF EVIDENCE TABLE** — Markdown table: [Study/Source, Design, Sample Size, Key Finding, Limitation]\n"
        "  3) **PRACTICAL SIGNIFICANCE** — What this means in the real world, not just statistically\n"
        "  4) **CRITICAL UNKNOWNS** — Top 3 gaps that would change the verdict if resolved\n"
        "  5) **RECOMMENDED STUDY DESIGNS** — Specific methodologies (RCT, cohort, meta-analysis) to close each gap\n"
        "STRICT RULE: Distinguish statistical significance from practical significance. Cite study designs by name."
    ) }
}

STARTUP_PHASE_DIRECTIVES = {
    0: { "title": "MARKET INTAKE", "instruction": "Size the TAM/SOM and map the competitive landscape with USD figures." },
    1: { "title": "UNIT ECONOMICS", "instruction": "Model LTV:CAC, burn rate, and funding requirements. Assume nothing." },
    2: { "title": "KILL THE BUSINESS", "instruction": "Find the fatal flaw: Market, Execution, Competition, or Timing risk." },
    3: { "title": "GO-TO-MARKET", "instruction": "Design the launch strategy: Beachhead segment, channels, and pricing." },
    4: { "title": "INVESTOR DECISION MEMO", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the LEAD PARTNER. Deliver the complete investment decision memo. Structure as:\n"
        "  1) **VERDICT** — Invest / Pass / Conditional (bold) with one sentence rationale\n"
        "  2) **VALUATION TABLE** — Markdown table: [Method, Pre-Money, Post-Money, Key Assumption]\n"
        "  3) **DEAL TERMS** — Recommended check size, instrument (SAFE/priced round), and key protective provisions\n"
        "  4) **FATAL FLAW RESPONSE** — The top risk identified and what evidence would de-risk it\n"
        "  5) **CONDITIONS TO CLOSE** — Specific milestones or diligence items required before wire\n"
        "STRICT RULE: Put a number on everything. No 'promising' or 'exciting' language."
    ) }
}

AUDIT_PHASE_DIRECTIVES = {
    0: { "title": "SCOPE & INVENTORY", "instruction": "Define audit universe and required evidence for SOC2/ISO/NIST." },
    1: { "title": "CONTROL TESTING", "instruction": "Test each control vs evidence. Findings: Effective / Ineffective / Exception." },
    2: { "title": "RED TEAM AUDIT", "instruction": "Attack the audit chain. Find self-reported evidence gaps." },
    3: { "title": "REMEDIATION", "instruction": "Map every finding to a specific fix, owner, and timeline." },
    4: { "title": "AUDIT OPINION — Full Report", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the LEAD AUDITOR. Deliver the complete audit opinion. Structure as:\n"
        "  1) **OPINION** — Clean / Qualified / Adverse / Disclaimer of Opinion (bold, one line)\n"
        "  2) **FINDINGS TABLE** — Markdown table: [Control, Status, Severity (Critical/High/Med/Low), Evidence Gap, Owner]\n"
        "  3) **MANAGEMENT LETTER POINTS** — Specific deficiencies that fall short of a finding but need attention\n"
        "  4) **REMEDIATION ROADMAP** — Numbered list: each finding with specific fix, deadline, and responsible party\n"
        "  5) **RE-AUDIT CRITERIA** — What must be true for a clean opinion on the next audit cycle\n"
        "STRICT RULE: Every finding must map to a specific control framework (SOC2 CC, ISO Annex A, NIST CSF). No vague conclusions."
    ) }
}

CREATIVE_PHASE_DIRECTIVES = {
    0: { "title": "STRATEGIC INSIGHT", "instruction": "Audience psychographics and the 'tension' that makes this relevant." },
    1: { "title": "CONCEPT DEVELOPMENT", "instruction": "Generate 3 Big Ideas with headlines, visuals, and tone sketches." },
    2: { "title": "CREATIVE CRITIQUE", "instruction": "Quality control: Originality check and Strategic fit audit." },
    3: { "title": "MAIN DRAFT", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "WRITE the actual finished deliverable (LinkedIn post, email, ad copy). "
        "Write as the brand/persona, match platform conventions, and sound like a human builder."
    )},
    4: { "title": "CREATIVE EXECUTION PACKAGE", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the CREATIVE DIRECTOR. Deliver the complete execution package. Structure as:\n"
        "  1) **SELECTED CONCEPT** — The winning Big Idea from phase 1, stated in one line\n"
        "  2) **FINAL DELIVERABLE** — The actual finished creative work (copy, script, campaign brief) ready to use\n"
        "  3) **ALTERNATE VERSION** — A second ready-to-use version with a different angle or tone\n"
        "  4) **ROLLOUT TIMELINE** — Markdown table: [Phase, Action, Owner, Date, Success KPI]\n"
        "  5) **MEASUREMENT PLAN** — Specific metrics and benchmarks to evaluate creative performance\n"
        "STRICT RULE: The deliverable must be ready to hand to a client or post immediately. No meta-commentary."
    ) }
}

TECH_PHASE_DIRECTIVES = {
    0: { "title": "ARCHITECTURE INTAKE", "instruction": "Map system components, data flows, and technical debt. Be specific: name technologies, versions, protocols. No hand-waving." },
    1: { "title": "TRADE-OFF ANALYSIS", "instruction": "Challenge the default tech choice. For every recommendation, state what you LOSE by choosing it. Name the cheaper, faster, or simpler alternative and explain why it was rejected. If the previous phase recommended something, find the weakness in that recommendation. No unanimous agreement — trade-offs always exist." },
    2: { "title": "FAILURE MODES", "instruction": "Your job is to BREAK the proposed architecture. You are the adversary. Find the single point of failure that will cause a 3AM outage. Identify the bottleneck that collapses under 10x load. Name the dependency that will be deprecated in 18 months. Challenge any claim from prior phases that lacks evidence. If Phase 1 said 'scalable', prove where it stops scaling. Be specific: name the component, the failure mode, and the blast radius. Do NOT agree with prior phases unless you have independent evidence." },
    3: { "title": "BUILD SEQUENCE", "instruction": "Phased implementation: MVP, Scale, and Harden cycles. Address every failure mode from Phase 2 — if you skip one, justify why." },
    4: { "title": "CTO DECISION BRIEF", "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

        "You are the CTO. Deliver the complete technical decision brief. Structure as:\n"
        "  1) **VERDICT** — Ship / Rearchitect / Stop (bold) with one-sentence rationale\n"
        "  2) **ARCHITECTURE DECISION TABLE** — Markdown table: [Component, Current State, Recommended Change, Effort (S/M/L), Risk]\n"
        "  3) **TECH DEBT REGISTER** — Accepted debt items with a specific payoff timeline and cost estimate\n"
        "  4) **BUILD SEQUENCE** — Ordered list: MVP → Scale → Harden with specific milestones\n"
        "  5) **GO/NO-GO CHECKLIST** — The exact criteria that must be met before shipping to production\n"
        "STRICT RULE: Name specific technologies, version numbers, and cost estimates. No generic 'refactor the codebase' language."
    ) }
}

SYSTEM_PHASE_DIRECTIVES = {
    0: { "title": "SYSTEM STATUS", "instruction": "Extract symptoms, affected flows, and monitoring metrics." },
    1: { "title": "ROOT CAUSE", "instruction": "Hypothesis ranking based on evidence and contradictory data." },
    2: { "title": "BLAST RADIUS", "instruction": "Scope downstream dependencies and hidden data corruption." },
    3: { "title": "RESOLUTION", "instruction": "Immediate fix actions, verification steps, and rollback plan." },
    4: { "title": "PREVENTION", "instruction": "Prevent recurrence: SIEM updates, arch changes, and runbook audits." }
}

EOM_STATEMENT_PHASE_DIRECTIVES = {
    0: {
        "title": "INTAKE — Financial Extraction",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the FORENSIC ACCOUNTANT. Your job is to extract every piece of raw financial data from the query. "
            "Build structured tables for Revenue, COGS, and OPEX. Identify every specific transaction or summary figure. "
            "Do NOT analyze yet. Just build the absolute ground-truth data table."
        )
    },
    1: {
        "title": "ANALYSIS — P&L & Balance Sheet Modeling",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the FINANCIAL ANALYST. Using the intake data, build a structured Profit & Loss statement. "
            "Include Gross Margin, EBITDA, and Net Income calculations. "
            "If balance sheet data (Assets/Liabilities) is present, build a snapshot. State all formulas used."
        )
    },
    2: {
        "title": "CASH FLOW — Burn & Runway Audit",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the CFO. Calculate the monthly burn rate (Gross and Net). "
            "Based on cash-on-hand, determine the exact runway in months. "
            "Identify the 'Zero Cash Date' and flag any immediate liquidity risks."
        )
    },
    3: {
        "title": "VARIANCE & STRATEGIC INSIGHT",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the STRATEGIC ADVISOR. Identify anomalies (spikes in OPEX, drops in margin). "
            "What happened this month that was unexpected? Provide 3 specific strategic levers to improve the next month's outcome."
        )
    },
    4: {
        "title": "VERDICT — Executive Financial Statement",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the CHIEF FINANCIAL OFFICER. Output the final, ready-to-present EOM report. "
            "Structure as: 1) Executive Summary (Financial Health Score), 2) Standard P&L Table, 3) Burn & Runway Metrics, 4) Top 3 Actions for next month. "
            "STRICT RULE: Be cold, analytical, and precise. No marketing language."
        )
    }
}

SOCIAL_POST_PHASE_DIRECTIVES = {
    0: {
        "title": "NARRATIVE HOOK — The Builder's Tension",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the NARRATIVE STRATEGIST. Your job is to find the 'tension' or 'contrarian insight' in the query. Output:\n"
            "  1) THE TENSION — the specific problem, irony, or hard truth that makes this post matter to builders/founders\n"
            "  2) THE ANGLE — a specific, non-obvious take (e.g. 'Why our biggest security feature is actually a design choice')\n"
            "  3) TARGET READER — be specific: Senior Devs, CISOs, Bootstrapped Founders, etc.\n"
            "Do NOT write the post. Find the hook that stops a builder from scrolling."
        )
    },
    1: {
        "title": "VOICE CALIBRATION — Zero Fluff Audit",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the VOICE ANALYST. Strip away all 'AI-voice' or 'marketer-voice'. Output:\n"
            "  1) FORBIDDEN PHRASES — list 5 generic buzzwords to avoid (e.g. 'game-changer', 'delve', 'unleash')\n"
            "  2) VOICE PROFILE — set the tone as 'Technical & Bare-metal' or 'Humble Founder' or 'Opinionated Expert'\n"
            "  3) AUTHENTICITY ANCHOR — one personal or technical detail from the query that proves this is written by a human\n"
            "Do NOT write the post. Set the constraints for a zero-fluff draft."
        )
    },
    2: {
        "title": "DRAFTING — The Builder's Narrative",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the CONTENT WRITER. Write three distinct, ready-to-publish versions of the post based on the hook and voice above.\n\n"
            "For EACH version output:\n"
            "  **VERSION [1/2/3]: [Tagline — e.g. 'The Raw Truth', 'The Technical Deep-Dive', 'The Contrarian Take']**\n"
            "  - The complete post text, formatted for the target platform (LinkedIn/X)\n"
            "  - Rule: Start with a 1-sentence hook. Use line breaks. No emojis unless essential.\n\n"
            "RULES:\n"
            "  - Sound like you've actually built something. Reference real trade-offs and failures.\n"
            "  - No generic overviews. Get specific immediately."
        )
    },
    3: {
        "title": "EDITOR — Cutting the Fat",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the LEAD EDITOR. Rank the versions and sharpen the winner.\n"
            "  1) RANKING — 1/2/3 with a 1-sentence reason (focus on authenticity)\n"
            "  2) SHARPENED WINNER — take the top post and cut 20% of the words. Make every line hit harder.\n"
            "  3) CTA — a natural, non-pushy ending that invites discussion, not clicks."
        )
    },
    4: {
        "title": "FINAL DELIVERY — Ready to Ship",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the DELIVERY LEAD. Output the final package for immediate posting:\n"
            "  1) **RECOMMENDED POST** — the sharpened winning version, ready to copy-paste\n"
            "  2) **ALTERNATE VERSION A** — the second-best version\n"
            "  3) **ALTERNATE VERSION B** — the third-best version\n"
            "  4) **POST METADATA** — target audience, best time to post, and recommended hook index.\n\n"
            "STRICT RULE: No meta-commentary. No 'Here are your posts'. Just the content. The 'RECOMMENDED POST' section must start with the actual first line of the post."
        )
    }
}

PORTFOLIO_BUILDER_PHASE_DIRECTIVES = {
    0: {
        "title": "MACRO SETUP — Market Environment",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the MACRO STRATEGIST. Read the current market environment cold. Output:\n"
            "  1) **MACRO SIGNAL** — Risk-on or Risk-off? One word, then one sentence rationale\n"
            "  2) **SECTOR ROTATION TABLE** — Markdown table: [Sector, Momentum (Hot/Neutral/Cold), Macro Tailwind/Headwind]\n"
            "  3) **RATE & INFLATION POSTURE** — Current Fed stance and what it means for equities vs bonds vs commodities\n"
            "  4) **TOP 3 MACRO RISKS** — The events that would blow up any portfolio in the next 90 days\n"
            "Do NOT pick stocks yet. Set the battlefield."
        )
    },
    1: {
        "title": "STOCK SCREENER — Candidate Selection",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the QUANT ANALYST. Screen and surface specific investment candidates based on the user's query and macro setup. Output:\n"
            "  1) **EQUITY CANDIDATES TABLE** — Markdown table: [Ticker, Company, Sector, P/E, Revenue Growth %, Moat, Catalyst]\n"
            "  2) **ETF/INDEX CANDIDATES** — If diversification is warranted, name specific ETF tickers with expense ratios\n"
            "  3) **CONTRARIAN PICKS** — 1-2 out-of-consensus names with a specific reason they're undervalued\n"
            "STRICT RULE: Real tickers only. No 'companies like X' — name the actual ticker."
        )
    },
    2: {
        "title": "RED TEAM — Kill the Picks",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the SHORT SELLER. Your job is to find the fatal flaw in every candidate from the screener. Output:\n"
            "  1) **BEAR CASE TABLE** — Markdown table: [Ticker, Bear Case Thesis, Probability (H/M/L), Kill Switch Event]\n"
            "  2) **VALUATION TRAPS** — Which picks look cheap but aren't? Name them and explain why\n"
            "  3) **SECTOR LANDMINES** — Upcoming earnings, macro events, or regulatory moves that could crater a position\n"
            "  4) **SURVIVORS** — After the red team, which picks still stand? List only those with a survivable bear case."
        )
    },
    3: {
        "title": "PORTFOLIO ARCHITECTURE — Position Sizing",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the RISK MANAGER. Design the portfolio structure around the surviving picks. Output:\n"
            "  1) **ALLOCATION TABLE** — Markdown table: [Ticker, Allocation %, Rationale, Max Drawdown Tolerance]\n"
            "  2) **CORRELATION MAP** — Flag any picks that move together (correlated risk) and adjust weightings\n"
            "  3) **HEDGES** — If the macro setup warrants it, name a specific hedge (put, inverse ETF, commodity) with sizing\n"
            "  4) **REBALANCE TRIGGERS** — The specific price or event that forces a position review\n"
            "STRICT RULE: Allocations must sum to 100%. Show your math."
        )
    },
    4: {
        "title": "INVESTMENT COMMITTEE VERDICT — Final Portfolio",
        "instruction": (
            "STRATEGIC RULE: Use [SIGNAL TAGS] (e.g. [CRITICAL], [VERIFIED], [ACTION REQUIRED]) to highlight key findings. "
            "Default to Amber [TAGS]. Use Red [CRITICAL] only for high-threat signals. "

            "You are the CHIEF INVESTMENT OFFICER. Compile the complete, ready-to-execute portfolio. Structure as:\n"
            "  1) **PORTFOLIO SUMMARY** — Total positions, cash %, expected 12-month return range, max drawdown estimate\n"
            "  2) **FINAL PORTFOLIO TABLE** — Markdown table: [Ticker, Name, Allocation %, Entry Price Range, 12-Month Target, Stop Loss, Conviction (1-10)]\n"
            "  3) **THESIS IN ONE LINE** — For each position, one sentence on why it wins\n"
            "  4) **EXIT STRATEGY** — For each position, the specific event or price that triggers a full exit\n"
            "  5) **WATCHLIST** — 2-3 names that didn't make the cut but should be monitored with entry conditions\n"
            "STRICT RULE: Every number must come from prior phase analysis. No generic diversification advice. "
            "Give the portfolio like you have skin in the game."
        )
    }
}

WORKFLOW_PHASE_OVERRIDES = {
    "FINANCE": FINANCE_PHASE_DIRECTIVES,
    "WAR_ROOM": WAR_ROOM_PHASE_DIRECTIVES,
    "LEGAL": LEGAL_PHASE_DIRECTIVES,
    "MEDICAL": MEDICAL_PHASE_DIRECTIVES,
    "QUANTUM_SECURITY": QUANTUM_SECURITY_PHASE_DIRECTIVES,
    "DEFENSE": DEFENSE_PHASE_DIRECTIVES,
    "CYBER": CYBER_PHASE_DIRECTIVES,
    "INTEL": INTEL_PHASE_DIRECTIVES,
    "SCIENCE": SCIENCE_PHASE_DIRECTIVES,
    "STARTUP": STARTUP_PHASE_DIRECTIVES,
    "AUDIT": AUDIT_PHASE_DIRECTIVES,
    "CREATIVE": CREATIVE_PHASE_DIRECTIVES,
    "TECH": TECH_PHASE_DIRECTIVES,
    "SYSTEM": SYSTEM_PHASE_DIRECTIVES,
    # --- ALIASES TO BRIDGE UI NAMES ---
    "CREATIVE_COUNCIL": CREATIVE_PHASE_DIRECTIVES,
    "STARTUP_LAUNCH": STARTUP_PHASE_DIRECTIVES,
    "CODE_AUDIT": AUDIT_PHASE_DIRECTIVES,
    "CYBER_COMMAND": CYBER_PHASE_DIRECTIVES,
    "DEFENSE_COUNCIL": DEFENSE_PHASE_DIRECTIVES,
    "INTEL_BRIEF": INTEL_PHASE_DIRECTIVES,
    "SCIENCE_PANEL": SCIENCE_PHASE_DIRECTIVES,
    "SOCIAL_POST": SOCIAL_POST_PHASE_DIRECTIVES,
    "EOM_STATEMENT": EOM_STATEMENT_PHASE_DIRECTIVES,
    "PORTFOLIO_BUILDER": PORTFOLIO_BUILDER_PHASE_DIRECTIVES,
}

class CouncilContext:
    def __init__(self, query, classification, workflow="RESEARCH", session_id=None, run_id=None, previous_context=None, user_id=None, ghost_map=None, residual_report=None):
        self.query = query
        self.classification = classification
        self.workflow = workflow.upper() if workflow else "RESEARCH"
        self.session_id = session_id
        self.run_id = run_id
        self.user_id = user_id
        self.history = []
        self.previous_context = previous_context or []
        # Falcon Ghost Map + PII Diff — injected when analyzing redacted documents
        self.ghost_map: dict = ghost_map or {}       # from build_ghost_map_summary()
        self.residual_report: dict = residual_report or {}  # from detect_residual_pii()

    def add_entry(self, ai_name, persona, response, usage=None):
        self.history.append({
            "ai": ai_name,
            "persona": persona,
            "response": response,
            "usage": usage or {},
            "timestamp": datetime.now().isoformat()
        })

def classify_query_deterministic(workflow, active_models=None):
    """Returns a locked, deterministic execution order for a workflow."""
    if not active_models:
        active_models = ["openai", "anthropic", "google", "perplexity", "mistral"]
    
    base_order = WORKFLOW_STEPS.get(workflow, WORKFLOW_STEPS["RESEARCH"])
    locked_order = []
    
    # Map requested providers to active ones
    for step in base_order:
        provider, role = step.split('-', 1)
        if provider in active_models:
            locked_order.append(step)
        else:
            # Fallback logic: if a primary model is missing, use another available model to fill the role.
            # This maintains the 5-phase structure even with limited models.
            # Use available models in rotation to ensure diversity where possible.
            fallback_idx = base_order.index(step) % len(active_models)
            fallback_provider = active_models[fallback_idx]
            locked_order.append(f"{fallback_provider}-{role}")
            
    return locked_order

# --- PHASE 1: CLASSIFICATION (The Planner) ---
def classify_query_v2(query, active_personas, active_models=None, previous_context=None, user_id=None, workflow="RESEARCH"):
    if active_models is None:
        active_models = ["openai", "anthropic", "google", "perplexity", "mistral"]

    # --- PRODUCTION LOCK: Deterministic Path ---
    # If the workflow is explicitly identified, return the locked pipeline immediately.
    # This bypasses LLM planning for all standard workflows, ensuring perfect consistency.
    if workflow in WORKFLOW_STEPS:
        print(f"[PLANNER] LOCK ACTIVE: Using deterministic pipeline for {workflow}")
        return {
            "domain": workflow.lower(),
            "intent": "analyze",
            "complexity": "complex",
            "outputType": "report",
            "executionOrder": classify_query_deterministic(workflow, active_models),
            "reasoning": f"Locked deterministic pipeline for production stability ({workflow})."
        }

    # --- DYNAMIC PATH: LLM-Based Classification (Fallback) ---
    available_list = ""
    for p in active_models:
        available_list += f"\n    - {p.capitalize()}: {active_personas.get(p, 'Analyst')}"

    # Build prior session context block if this is a follow-up query
    prior_block = ""
    if previous_context:
        prior_block = "\n    FOLLOW-UP CONTEXT — This query builds on a previous council session:\n"
        for entry in previous_context[-2:]:
            prior_block += f"    - Previous Query: \"{entry.get('query', 'N/A')}\"\n"
            prior_block += f"    - Consensus Score: {entry.get('consensus_score', 'N/A')}/100\n"
            prior_block += f"    - Summary: {entry.get('summary', 'N/A')}\n"
            contested = entry.get('contested_topics', [])
            if contested:
                prior_block += f"    - Unresolved Topics: {', '.join(contested)}\n"
        prior_block += "    The council should build on these prior conclusions, not repeat them.\n"

    prompt = f"""
    Analyze this query and determine optimal AI execution order.

    QUERY: "{query}"
{prior_block}

    AVAILABLE PERSONAS (User Selected):{available_list}

    OPTIMAL EXECUTION ORDER PRINCIPLES (STRICT 5-PHASE PIPELINE):
    Phase 1 - INTAKE: Neutral baseline analysis — best for a broad strategist or analyst (prefer OpenAI)
    Phase 2 - STRATEGIC INTERPRETATION: Build scenarios from intake — best for a deep analyst or architect (prefer Anthropic)
    Phase 3 - CHALLENGE: Stress-test assumptions from Phase 2 — best for a critic, researcher, or scout (prefer Google or Perplexity)
    Phase 4 - OPERATIONS: Translate analysis into actionable steps — best for an operator, scout, or domain specialist (prefer Perplexity or Mistral)
    Phase 5 - VALIDATION: Final quality check, framework mapping, confidence assessment — best for a critic or validator (prefer Mistral)

    ROLE SELECTION RULES:
    - Match each persona's role to the DOMAIN of the query, not to a fixed specialty.
    - A medical query should use roles like medical, researcher, analyst — NOT cryptographer or zero_trust.
    - A finance query should use roles like cfo, auditor, economist — NOT cyber_ops or defense_ops.
    - Only assign security/crypto roles (cryptographer, zero_trust, cyber_ops, hacker) when the query is ACTUALLY about cybersecurity or cryptography.
    - Maximize DIVERSITY across the council — avoid assigning similar roles to multiple providers.

    IMPORTANT: Only use the personas provided in the AVAILABLE PERSONAS list. Do not use any others. Include all available personas in the executionOrder.

    return ONLY valid JSON (no markdown):
    {{
      "domain": "business|marketing|software|operations|research|strategy|medical|legal|creative|engineering|science|finance|cybersecurity",
      "intent": "plan|build|analyze|optimize|critique|research|launch|design|assess",
      "complexity": "simple|moderate|complex",
      "outputType": "presentation|technical_spec|marketing_plan|report|strategic_framework|diagram|creative_brief|research_paper",
      "executionOrder": ["provider-role", "provider-role"],
      "reasoning": "Brief explanation of order and why these roles match the query domain"
    }}
    """
    
    # 1. Try OpenAI (Primary)
    try:
        plan_response = call_openai_gpt4(prompt, "Planner", user_id=user_id, timeout=20)
        if plan_response['success']:
             content = plan_response['response'].replace('```json', '').replace('```', '').strip()
             return json.loads(content)
        else:
             print("[PLANNER] OpenAI Failed. Trying Mistral API...")
    except Exception as e:
        print(f"[PLANNER ERROR] OpenAI: {e}")

    # 2. Try Mistral API (Cloud Fallback)
    try:
        plan_response = call_mistral_api(prompt, "Planner", user_id=user_id, timeout=20)
        if plan_response['success']:
             content = plan_response['response'].replace('```json', '').replace('```', '').strip()
             return json.loads(content)
        else:
             print("[PLANNER] Mistral API Failed. Trying Local...")
    except Exception as e:
        print(f"[PLANNER ERROR] Mistral API: {e}")

    # 3. Try Local LLM (Emergency Fallback)
    try:
        local_response = call_local_llm(prompt, "Planner", user_id=user_id)
        if local_response['success']:
            content = local_response['response'].replace('```json', '').replace('```', '').strip()
            # Local models can be chatty, try to find JSON blob
            start = content.find('{')
            end = content.rfind('}') + 1
            if start != -1 and end != -1:
                return json.loads(content[start:end])
    except Exception as local_e:
        print(f"[PLANNER ERROR] Local: {local_e}")

    # Ultimate Fallback - Use at least one cloud provider (Mistral) as a scout if possible
    return {
        "executionOrder": ["mistral-scout", "mistral-strategist", "mistral-critic"],
        "reasoning": "PRIMARY MODELS OFFLINE. Using Mixed Emergency Council.",
        "outputType": "report"
    }


# --- PHASE 2: SEQUENTIAL EXECUTION (The Runner) ---
# --- PHASE 8: AI ACCOUNTABILITY (The Enforcer) ---

def identify_claims(text, ai_name, user_id=None):
    """
    Parses AI response for specific, testable claims.
    Uses a lightweight LLM call or regex to find metrics/absolute statements.
    """
    prompt = f"""
    Identify specific, testable claims from this AI response ({ai_name}).
    Focus on:
    1. Numerical metrics (e.g. "40% growth")
    2. Absolute statements (e.g. "This is the ONLY way...")
    3. Proper nouns/entities
    4. Causal links (e.g. "X causes Y")

    RESPONSE:
    {text}

    Return ONLY a JSON list of claims:
    [
      {{"claim": "Claim text", "type": "metric|absolute|entity|causal"}}
    ]
    """
    try:
        # Use GPT-4o-mini for speed/consistency in parsing
        resp = call_openai_gpt4(prompt, "ClaimExtractor", model="gpt-4o-mini", user_id=user_id)
        if resp['success']:
            content = resp['response'].replace('```json', '').replace('```', '').strip()
            return json.loads(content)
    except Exception as e:
        print(f"[CLAIM ERROR] {e}")
    return []

def verify_claims(claims, council_history):
    """
    Cross-references claims against other advisor outputs.
    Labels them [CONFIRMED], [SUSPECT], or [UNVERIFIED].

    Each claim gets:
    - 'contribution': +/- impact on truth score
    - 'confidence': HIGH / MEDIUM / LOW per-claim label
    - 'lifecycle': original → challenged → verified → updated
    - 'violations': list of enforcement flags (PQC, probability, temporal)
    """
    verified_results = []

    # Regex patterns for probability and temporal enforcement
    _prob_pattern = re.compile(r'(\d{1,3})\s*%')
    _prob_basis_signals = [
        'based on', 'according to', 'data shows', 'model predicts', 'analysis indicates',
        'historical', 'calculated', 'derived from', 'regression', 'forecast',
        'survey', 'sample', 'study', 'report', 'source:', 'per '
    ]
    _temporal_vague = re.compile(
        r'\b(recently|recent|increasing|increasingly|growing|declining|'
        r'trending|surging|emerging|rising|falling|accelerating|'
        r'more and more|on the rise|gaining traction)\b', re.IGNORECASE
    )
    _temporal_grounded = re.compile(
        r'\b(20\d{2}|Q[1-4]|january|february|march|april|may|june|'
        r'july|august|september|october|november|december|'
        r'year-over-year|YoY|month-over-month|MoM|since \d{4}|'
        r'last \d+ (years?|months?|quarters?|weeks?))\b', re.IGNORECASE
    )

    for c in claims:
        claim_text = c['claim']
        claim_lower = claim_text.lower()
        status = "UNVERIFIED"
        score = 50
        anchors = []
        violations = []
        contribution = 0  # net +/- impact on truth score

        # --- QUANTUM DRIFT CHECK (FIPS 203/204/205/206) ---
        legacy_crypto = ["RSA", "ECC", "ECDSA", "Diffie-Hellman", "AES-128", "DSA", "3DES", "RC4", "MD5", "SHA-1"]
        pqc_wrappers = ["ML-KEM", "Kyber", "ML-DSA", "Dilithium", "SLH-DSA", "Sphincs+", "FALCON", "FN-DSA"]

        has_legacy = any(lc.lower() in claim_lower for lc in legacy_crypto)
        has_pqc = any(pqc.lower() in claim_lower for pqc in pqc_wrappers)

        if has_legacy and not has_pqc:
            violations.append("Non-PQC Compliant: Legacy cryptography detected without FIPS 203/204/205/206 wrapper (Quantum Drift). Recommend ML-KEM (key encapsulation), SLH-DSA (long-term signing), or FALCON/FN-DSA as Integrity Anchor for constrained-bandwidth paths.")
            score -= 20
            contribution -= 8  # PQC violation is a significant drag

        # --- PROBABILITY VALIDATION (Item 2) ---
        # Any percentage claim must have a source, model, or reasoning basis
        prob_matches = _prob_pattern.findall(claim_text)
        if prob_matches:
            has_basis = any(sig in claim_lower for sig in _prob_basis_signals)
            if not has_basis:
                violations.append(f"UNSUPPORTED PROBABILITY: '{prob_matches[0]}%' stated without source, model, or reasoning basis. Auto-downgraded.")
                score -= 10
                contribution -= 4

        # --- TEMPORAL CLAIM ENFORCEMENT (Item 3) ---
        # Block vague time-based language without evidence
        temporal_vague_hits = _temporal_vague.findall(claim_text)
        if temporal_vague_hits:
            has_temporal_ground = bool(_temporal_grounded.search(claim_text))
            if not has_temporal_ground:
                violations.append(f"UNGROUNDED TEMPORAL CLAIM: '{temporal_vague_hits[0]}' used without time-bound evidence or data reference.")
                score -= 5
                contribution -= 2

        # Cross-provider agreement via keyword overlap
        # Extract significant words (4+ chars, not stopwords) from the claim
        _stopwords = {'this', 'that', 'with', 'from', 'will', 'have', 'been', 'they',
                      'their', 'would', 'could', 'should', 'about', 'which', 'there',
                      'these', 'those', 'than', 'then', 'also', 'into', 'more', 'most',
                      'such', 'when', 'what', 'some', 'only', 'very', 'just', 'over',
                      'each', 'does', 'were', 'being', 'other', 'while', 'both', 'after',
                      'before', 'between', 'under', 'through', 'during', 'without', 'within'}
        claim_words = [w for w in re.findall(r'[a-z0-9]+', claim_lower)
                       if len(w) >= 4 and w not in _stopwords]
        keyword_threshold = max(2, len(claim_words) * 0.4)  # 40% of keywords must match

        agreement_count = 0
        for entry in council_history:
            content = entry['response'].lower()
            matched = sum(1 for w in claim_words if w in content)
            if claim_words and matched >= keyword_threshold:
                agreement_count += 1
                anchors.append(entry['ai'])

        # Scoring Logic — with claim-level contribution tracking
        if agreement_count >= 2:
            status = "CONFIRMED"
            score += 40
            contribution += 5   # confirmed claims boost truth
        elif agreement_count == 1:
            status = "SUPPORTED"
            score += 25
            contribution += 2   # partial support — modest boost
        else:
            # UNVERIFIED — no cross-provider support
            contribution -= 3   # unverified claims drag truth down

        # Type-based weight adjustments
        claim_type = c.get('type', 'unknown')
        if claim_type == 'metric' and status == 'UNVERIFIED':
            contribution -= 2   # unverified numbers are high risk
        elif claim_type == 'absolute' and status != 'CONFIRMED':
            contribution -= 3   # absolute statements need strong backing
        elif claim_type == 'causal' and status == 'CONFIRMED':
            contribution += 2   # confirmed causal claims add high value

        # --- CONFIDENCE LABELING (Item 7) ---
        # Per-claim confidence based on final score
        final_score = max(0, min(score, 100))
        if final_score >= 75:
            confidence = "HIGH"
        elif final_score >= 50:
            confidence = "MEDIUM"
        else:
            confidence = "LOW"

        # --- CLAIM LIFECYCLE (Item 5) ---
        # Track where this claim is in its lifecycle
        # Default is 'original' — other states set by mediation/interrogation
        lifecycle = "original"
        if violations:
            lifecycle = "challenged"

        verified_results.append({
            "claim": claim_text,
            "status": status,
            "score": final_score,
            "type": claim_type,
            "anchors": anchors,
            "violations": violations,
            "contribution": contribution,
            "confidence": confidence,
            "lifecycle": lifecycle
        })

    return verified_results


def calculate_truth_score(verified_claims):
    """
    Calculates overall card confidence using claim-level contributions.

    Returns int (0-100) truth score. The score is an aggregate:
    - Baseline: 70 (neutral starting point — not guilty, not proven)
    - Each claim's contribution is added/subtracted
    - Clamped to 0-100 range

    Also enriches each claim dict with 'contribution_pct' for display.
    """
    if not verified_claims:
        return 70  # Neutral baseline — no claims to evaluate

    baseline = 70
    total_contribution = sum(c.get('contribution', 0) for c in verified_claims)
    raw_score = baseline + total_contribution
    final_score = max(0, min(100, raw_score))

    # Enrich claims with contribution percentage for frontend display
    for c in verified_claims:
        c['contribution_pct'] = round(
            (c.get('contribution', 0) / max(1, abs(total_contribution))) * 100
        ) if total_contribution != 0 else 0

    return final_score


# ===================================================================
# SCORE MEDIATION ENGINE — Bidirectional Truth Recalibration
# Scores are a living instrument: down on weakness, up on correction.
# ===================================================================

def mediate_truth_scores(current_results, prior_events):
    """
    Compare current council results against prior interrogation/verification
    events from the thread. Produce per-provider score deltas.

    prior_events: list of dicts with keys:
        - type: 'interrogation' | 'verification'
        - provider: provider key (e.g. 'openai')
        - score_delta: int (the original penalty/bonus applied)
        - verdict: str (e.g. 'CONCESSION', 'INACCURATE', etc.)
        - keywords: list[str] — key terms from the contested claim

    Returns dict: {
        provider: {
            'delta': int,           # net score change to apply
            'reason': str,          # human-readable explanation
            'components': list,     # breakdown of what moved
            'recovery_pct': float   # 0.0-1.0 if recovering
        }
    }
    """
    if not prior_events or not current_results:
        return {}

    mediation = {}

    # Group prior negative events by provider
    provider_penalties = {}
    for evt in prior_events:
        prov = evt.get('provider')
        delta = evt.get('score_delta', 0)
        if not prov or delta >= 0:
            continue  # only track prior penalties
        if prov not in provider_penalties:
            provider_penalties[prov] = []
        provider_penalties[prov].append(evt)

    for provider, penalties in provider_penalties.items():
        result = current_results.get(provider)
        if not result or not result.get('success') or not result.get('response'):
            continue

        response_lower = result['response'].lower()
        components = []
        total_recovery = 0

        for penalty in penalties:
            penalty_amount = abs(penalty['score_delta'])
            verdict = penalty.get('verdict', '')
            keywords = penalty.get('keywords', [])

            # Check if the current response addresses the prior weakness
            addressed = False
            address_strength = 0.0

            # Signal 1: Response contains correction language referencing the issue
            correction_signals = [
                'upon further analysis', 'revised assessment', 'updated finding',
                'corrected', 'additional evidence', 'further review',
                'refined analysis', 'upon reflection', 'taking into account',
                'more accurate', 'better supported', 'revised estimate',
                'additional data supports', 'evidence now indicates',
                'recalibrated', 'adjusted based on'
            ]
            correction_hits = sum(1 for s in correction_signals if s in response_lower)

            # Signal 2: Keywords from the contested claim appear (topic is addressed)
            keyword_hits = 0
            if keywords:
                keyword_hits = sum(1 for kw in keywords if kw.lower() in response_lower)
                keyword_coverage = keyword_hits / len(keywords) if keywords else 0
            else:
                keyword_coverage = 0.5  # no keywords tracked — give partial credit

            # Signal 3: Response shows hedged/calibrated language (learned from penalty)
            calibration_signals = [
                'moderate likelihood', 'insufficient data', 'limited evidence',
                'estimated range', 'preliminary', 'subject to', 'uncertain',
                'requires further', 'cannot confirm', 'approximate',
                'with caveats', 'contingent on'
            ]
            calibration_hits = sum(1 for s in calibration_signals if s in response_lower)

            # Determine if the weakness was addressed
            if correction_hits >= 2 and keyword_coverage >= 0.3:
                addressed = True
                address_strength = min(0.6, 0.3 + (correction_hits * 0.05) + (keyword_coverage * 0.15))
            elif correction_hits >= 1 and keyword_coverage >= 0.5:
                addressed = True
                address_strength = min(0.5, 0.25 + (keyword_coverage * 0.15))
            elif keyword_coverage >= 0.6 and calibration_hits >= 1:
                addressed = True
                address_strength = min(0.45, 0.2 + (calibration_hits * 0.05) + (keyword_coverage * 0.1))

            if addressed:
                recovery = int(penalty_amount * address_strength)
                recovery = max(1, recovery)  # at least 1 point recovery
                total_recovery += recovery
                components.append({
                    'type': 'RECOVERY',
                    'original_penalty': -penalty_amount,
                    'recovery': +recovery,
                    'recovery_pct': round(address_strength, 2),
                    'reason': f"Prior {verdict} ({-penalty_amount}) — corrected (+{recovery}, {int(address_strength*100)}% recovery)"
                })
            else:
                # Weakness not addressed — small ongoing drag
                drag = -1 if penalty_amount >= 8 else 0
                if drag:
                    total_recovery += drag
                    components.append({
                        'type': 'UNRESOLVED_DRAG',
                        'original_penalty': -penalty_amount,
                        'drag': drag,
                        'reason': f"Prior {verdict} ({-penalty_amount}) — unaddressed, ongoing drag ({drag})"
                    })

        # Check for NEW overconfidence in the follow-up (score can also go down)
        overconfidence_signals = [
            'guaranteed', 'absolutely certain', '100% probability',
            'without any doubt', 'impossible to fail', 'zero risk',
            'no chance of failure', 'certain to succeed'
        ]
        overconfidence_hits = sum(1 for s in overconfidence_signals if s in response_lower)
        if overconfidence_hits >= 1:
            oc_penalty = -(overconfidence_hits * 3)
            total_recovery += oc_penalty
            components.append({
                'type': 'OVERCONFIDENCE_PENALTY',
                'penalty': oc_penalty,
                'reason': f"New overconfident language detected ({oc_penalty})"
            })

        if components and total_recovery != 0:
            if total_recovery > 0:
                reason = f"Score recovery: +{total_recovery} (prior weakness addressed)"
            elif total_recovery < 0:
                reason = f"Score drag: {total_recovery} (unresolved issues persist)"
            else:
                reason = "Score unchanged (mixed signals)"

            mediation[provider] = {
                'delta': total_recovery,
                'reason': reason,
                'components': components,
                'recovery_pct': round(total_recovery / max(1, sum(abs(p['score_delta']) for p in penalties)), 2)
            }

    return mediation


# ===================================================================
# FINAL ARBITER ENGINE — Mathematical Consistency Validation
# Section 8a enforcement: deterministic, zero-LLM, regex-based.
# ===================================================================

def _resolve_canonical(label):
    """Map a raw label string to a canonical metric name via METRIC_ALIASES."""
    label_lower = label.lower().strip()
    # Require the label to be at least 2 chars to avoid noise
    if len(label_lower) < 2:
        return None
    for canonical, aliases in METRIC_ALIASES.items():
        for alias in aliases:
            if alias == label_lower or label_lower.endswith(alias) or alias.endswith(label_lower):
                return canonical
    return None


# Regex: captures metric labels followed by monetary/percentage values.
# Matches: "Revenue of $2.5M", "Net Profit: $700K", "ROI: 38.9%", "Total Costs — $1.8 million"
_METRIC_PATTERN = re.compile(
    r'(?P<label>[A-Za-z][A-Za-z\s\-/]{1,40}?)\s*'    # metric label (1-40 chars, starts with letter)
    r'[:\-—=]+\s*'                                      # separator (colon, dash, equals)
    r'(?P<sign>[-−])?'                                   # optional negative
    r'\s*(?P<currency>\$|USD\s?)?'                       # optional currency
    r'(?P<number>[\d,]+(?:\.\d+)?)'                      # the number (decimal only if digits follow dot)
    r'\s*(?P<scale>(?:billion|million|thousand|[KkMmBb])(?:n)?|%)?',  # scale/unit (word boundary safe)
    re.IGNORECASE
)

# Secondary pattern: "of $X", "is $X" style
_METRIC_PATTERN_OF = re.compile(
    r'(?P<label>[A-Za-z][A-Za-z\s\-/]{1,40}?)\s+'      # metric label
    r'(?:of|is|was|at|equals?|totals?|reached)\s+'       # connector word
    r'(?P<sign>[-−])?'
    r'\s*(?P<currency>\$|USD\s?)?'
    r'(?P<number>[\d,]+(?:\.\d+)?)'                      # the number (decimal only if digits follow dot)
    r'\s*(?P<scale>(?:billion|million|thousand|[KkMmBb])(?:n)?|%)?',  # scale/unit
    re.IGNORECASE
)


def extract_metrics(text, provider_name):
    """
    Extract numerical metric claims from provider text using regex.
    Zero LLM cost, deterministic. Returns list of extraction dicts.
    """
    extractions = []
    seen = set()  # dedup: (canonical, value) per provider

    for pattern in (_METRIC_PATTERN, _METRIC_PATTERN_OF):
        for match in pattern.finditer(text):
            label = match.group('label').strip()
            canonical = _resolve_canonical(label)
            if not canonical:
                continue

            try:
                raw_number = float(match.group('number').replace(',', ''))
            except ValueError:
                continue

            scale = (match.group('scale') or '').lower()

            # Apply scale multiplier
            if scale.startswith('k') or scale.startswith('thousand'):
                raw_number *= 1_000
            elif scale.startswith('m') or scale.startswith('million'):
                raw_number *= 1_000_000
            elif scale.startswith('b') or scale.startswith('billion'):
                raw_number *= 1_000_000_000
            elif scale.startswith('t') and not scale.startswith('thousand'):
                raw_number *= 1_000_000_000_000

            if match.group('sign') in ('-', '−'):
                raw_number = -raw_number

            unit = "%" if scale == "%" else ("$" if match.group('currency') else "units")

            # Dedup within same provider
            dedup_key = (canonical, round(raw_number, 2))
            if dedup_key in seen:
                continue
            seen.add(dedup_key)

            # Context snippet
            start = max(0, match.start() - 50)
            end = min(len(text), match.end() + 50)

            extractions.append({
                "provider": provider_name,
                "metric_name": canonical,
                "raw_label": label,
                "value": raw_number,
                "unit": unit,
                "context_snippet": text[start:end].strip(),
            })

    return extractions


def validate_math_identities(resolved_metrics):
    """
    Check arithmetic relationships between resolved metrics.
    Returns list of violation dicts. Zero LLM cost.
    """
    violations = []
    for (result_metric, op, operand_a, operand_b) in MATH_IDENTITIES:
        if result_metric not in resolved_metrics:
            continue
        if operand_a not in resolved_metrics or operand_b not in resolved_metrics:
            continue

        a = resolved_metrics[operand_a]
        b = resolved_metrics[operand_b]
        stated = resolved_metrics[result_metric]

        if op == "subtract":
            computed = a - b
        elif op == "divide_pct":
            computed = (a / b * 100) if b != 0 else None
        elif op == "add":
            computed = a + b
        else:
            continue

        if computed is None:
            continue

        # Tolerance: 2% of the expected value or $1000 (for dollar amounts), 2 points (for %)
        if op == "divide_pct":
            tolerance = max(abs(stated) * 0.02, 2.0)
        else:
            tolerance = max(abs(stated) * 0.02, 1000)

        delta = abs(computed - stated)
        if delta > tolerance:
            violations.append({
                "identity": f"{operand_a} {op} {operand_b} = {result_metric}",
                "expected": round(stated, 2),
                "computed": round(computed, 2),
                "delta": round(delta, 2),
                "corrected_values": {result_metric: round(computed, 2)},
            })

    return violations


def final_arbiter(results, divergence=None):
    """
    FINAL ARBITER ENGINE — Mathematical consistency validation layer.
    Runs AFTER divergence analysis, BEFORE synthesis.

    1. Extracts numerical metrics from all provider outputs (regex, no LLM)
    2. Detects cross-provider conflicts on the same metric (>5% spread)
    3. Validates arithmetic identities (revenue - costs = profit, etc.)
    4. Produces a resolution report + narrative directive for synthesis

    Returns ArbiterReport dict. Non-destructive — never mutates provider data.
    """
    all_extractions = []

    # Phase 1: Extract metrics from each provider
    for provider, result in results.items():
        if not isinstance(result, dict) or not result.get('success'):
            continue
        text = result.get('response', '')
        if not text:
            continue
        extractions = extract_metrics(text, provider)
        all_extractions.extend(extractions)

    if not all_extractions:
        return {
            "status": "CLEAN",
            "metric_inventory": {},
            "conflicts": [],
            "math_violations": [],
            "resolved_metrics": {},
            "narrative_directive": "",
        }

    # Phase 2: Group by canonical metric name
    metric_inventory = {}
    for ext in all_extractions:
        name = ext['metric_name']
        if name not in metric_inventory:
            metric_inventory[name] = []
        metric_inventory[name].append(ext)

    # Phase 3: Detect conflicts and resolve
    conflicts = []
    resolved_metrics = {}

    for metric_name, entries in metric_inventory.items():
        # First occurrence per provider
        provider_values = {}
        for e in entries:
            if e['provider'] not in provider_values:
                provider_values[e['provider']] = e['value']

        values = list(provider_values.values())
        if len(values) == 0:
            continue

        if len(values) == 1:
            resolved_metrics[metric_name] = values[0]
            continue

        spread = max(values) - min(values)
        mean_val = sum(values) / len(values)
        spread_pct = (spread / abs(mean_val) * 100) if mean_val != 0 else 0

        if spread_pct > 5:
            # Conflict: pick provider with highest truth_meter
            best_provider = max(
                provider_values.keys(),
                key=lambda p: results.get(p, {}).get('truth_meter', 0)
            )
            conflicts.append({
                "metric": metric_name,
                "providers": provider_values,
                "spread": round(spread, 2),
                "spread_pct": round(spread_pct, 1),
                "resolution": "HIGHEST_TRUTH",
                "resolved_value": provider_values[best_provider],
            })
            resolved_metrics[metric_name] = provider_values[best_provider]
        else:
            # Within tolerance — use mean
            resolved_metrics[metric_name] = round(mean_val, 2)

    # Phase 4: Validate arithmetic identities
    math_violations = validate_math_identities(resolved_metrics)

    # Override resolved values with computed (math wins)
    for v in math_violations:
        for metric_name, corrected_val in v['corrected_values'].items():
            resolved_metrics[metric_name] = corrected_val
            for c in conflicts:
                if c['metric'] == metric_name:
                    c['resolution'] = "MATH_DERIVED"
                    c['resolved_value'] = corrected_val

    # Phase 5: Build narrative directive for synthesis
    status = "CLEAN"
    directive_parts = []

    if conflicts:
        status = "CONFLICTS_DETECTED"
        directive_parts.append("⚠ NUMERICAL CONFLICTS DETECTED — use ONLY these resolved values:")
        for c in conflicts:
            prov_str = ", ".join(f"{p}={v}" for p, v in c['providers'].items())
            directive_parts.append(
                f"  • {c['metric'].upper()}: USE {c['resolved_value']} "
                f"(providers disagreed: {prov_str} — resolved via {c['resolution']})"
            )

    if math_violations:
        status = "MATH_VIOLATIONS"
        directive_parts.append("⚠ ARITHMETIC VIOLATIONS CORRECTED:")
        for v in math_violations:
            directive_parts.append(
                f"  • {v['identity']}: stated={v['expected']}, computed={v['computed']}. "
                f"USE the computed value."
            )

    if not conflicts and not math_violations:
        directive_parts.append("All numerical metrics are consistent across providers. No corrections needed.")

    return {
        "status": status,
        "metric_inventory": {k: len(v) for k, v in metric_inventory.items()},  # counts only (keep payload small)
        "conflicts": conflicts,
        "math_violations": math_violations,
        "resolved_metrics": resolved_metrics,
        "narrative_directive": "\n".join(directive_parts),
    }


# Falcon-token sanitizer — strips hallucinated redaction placeholders when Falcon is OFF.
# Matches patterns like [CURRENCY_AMOUNT_A1B2C3], [ORG_5D0176], [PERSON_622F9F], etc.
# Does NOT touch legitimate system tags ([METRIC_ANCHOR], [RISK_VECTOR], etc.).
_FAKE_FALCON_RE = re.compile(
    r'\[(?:CURRENCY_AMOUNT|ORG|PERSON|PROPER_NOUN|LOCATION|EMAIL|PHONE|SSN|DOB|'
    r'CREDIT_CARD|IP_ADDRESS|ALNUM_TAG|MEDICAL|LEGAL_ID)_[A-Fa-f0-9]{4,}\]'
)

def _strip_phantom_tokens(text):
    """Replace hallucinated Falcon-style tokens with '[value not provided]'."""
    return _FAKE_FALCON_RE.sub('[value not provided]', text)


def normalization_layer(text):
    """
    KorumOS Normalization Layer (Section 5 Directive)
    - Converts ALL markdown tables -> structured JSON objects
    - Removes raw markdown/text artifacts
    - Ensures render layer never receives raw AI output
    """
    if not text:
        return ""
        
    import re
    import json

    # 1. Convert markdown tables to structured JSON blocks
    # Pattern: | col | col | \n |---|---| \n | val | val |
    table_pattern = re.compile(r'(\|.*\|[\r\n]+\|[-:| ]+\|[\r\n]+(?:\|.*\|[\r\n]*)+)', re.MULTILINE)
    
    def table_replacer(match):
        table_text = match.group(0).strip()
        rows = [r.strip() for r in table_text.split('\n') if r.strip()]
        if len(rows) < 3: 
            return table_text
        
        # Extract headers and data
        header_row = [h.strip().replace('|', '') for h in rows[0].split('|') if h.strip()]
        data_rows = []
        for r in rows[2:]:
            cols = [c.strip() for c in r.split('|') if c.strip()]
            if len(cols) == len(header_row):
                data_rows.append(dict(zip(header_row, cols)))
            elif len(cols) > 0:
                # Pad or truncate if alignment is off
                padded = (cols + [""] * len(header_row))[:len(header_row)]
                data_rows.append(dict(zip(header_row, padded)))
        
        # Return structured block for frontend/exporter parsing
        return f"\n[STRUCTURED_TABLE]{json.dumps(data_rows)}[/STRUCTURED_TABLE]\n"

    # Run table normalization
    normalized = table_pattern.sub(table_replacer, text)
    
    # 2. Cleanup raw markdown leaks (Sections 6, 12)
    # Strip triple backticks except for specific allowed code blocks (if any)
    # We strip them to normalize the output into clean text.
    normalized = re.sub(r'```[a-z]*', '', normalized)
    normalized = normalized.replace('```', '')
    
    # 3. Strip bold/italic markers from report bodies for total uniform professionalism
    # (Optional, but often requested for "no markdown leakage")
    # Let's keep them for now but ensure structural tags are clean.
    
    return normalized.strip()


def execute_council_v2(query, active_personas, images=None, workflow="RESEARCH", active_models=None, previous_context=None, session_id=None, run_id=None, user_id=None, ledger_mission_id=None, ghost_map=None, residual_report=None):
    # 1. Setup IDs
    import uuid
    import hashlib as _hl
    if not run_id: run_id = str(uuid.uuid4())

    # Ledger helper — safe to call even if ledger is not active
    def _ledger_write(event_type, payload):
        if not ledger_mission_id:
            return
        try:
            from ledger import LedgerService
            LedgerService.write_event(
                event_type=event_type,
                mission_id=ledger_mission_id,
                decision_id=run_id,
                operator_id=user_id,
                payload=payload,
            )
        except Exception as e:
            print(f"[LEDGER] Warning: {event_type} write failed: {e}")

    # 2. Plan
    classification = classify_query_v2(query, active_personas, active_models=active_models, previous_context=previous_context, user_id=user_id, workflow=workflow)
    context = CouncilContext(query, classification, workflow=workflow, session_id=session_id, run_id=run_id, previous_context=previous_context, user_id=user_id, ghost_map=ghost_map, residual_report=residual_report)

    # --- LEDGER: mission_created (context established for council execution) ---
    _ledger_write("mission_created", {
        "mission_type": workflow,
        "classification": {
            "domain": classification.get("domain"),
            "intent": classification.get("intent"),
            "complexity": classification.get("complexity"),
            "output_type": classification.get("outputType"),
        },
    })

    if previous_context:
        print(f"[COUNCIL] Follow-up mode: {len(previous_context)} prior session(s) loaded")
    results = {}
    history_text = ""
    dna = {}
    outcome_goal = "Intelligence Brief"

    print(f"[COUNCIL] Efficiency Plan: {classification['executionOrder']}")
    if images:
        print(f"[COUNCIL] {len(images)} image(s) attached — vision mode active")

    total_run_cost = 0.0
    total_latency_ms = 0
    models_used = []

    # 3. Execute Step-by-Step
    try:
        execution_order = classification.get('executionOrder', [])
        if not isinstance(execution_order, list): execution_order = []

        # Safety net: backfill any providers the planner skipped
        if active_models is None:
            active_models = list(active_personas.keys())
        assigned_providers = set()
        for pr in execution_order:
            try:
                p = pr.split('-', 1)[0].lower().strip()
                assigned_providers.add(p)
            except:
                pass
        for model in active_models:
            if model.lower() == 'mistral':
                # Mistral is reserved for 'Shadow' roles in the council and the Finalizer pass.
                # We do NOT skip it here so it can be added to the execution_order if assigned a role.
                pass 
            if model.lower() not in assigned_providers:
                role = active_personas.get(model, 'analyst')
                execution_order.append(f"{model}-{role}")
                print(f"[COUNCIL] Backfilled missing provider: {model.upper()} as {role.upper()}")

        total_steps = len(execution_order)

        for i, provider_role in enumerate(execution_order):
            try:
                provider, role = provider_role.split('-', 1)
                provider = provider.lower().strip()
                role = role.strip()
            except ValueError:
                provider = provider_role.lower().strip()
                role = active_personas.get(provider, 'analyst')

            if i == total_steps - 1:
                role = f"Integrator ({role})"

            # --- SHADOW EXECUTION LOGIC (Section 10) ---
            # Mistral performs its assigned role but remains hidden from the UI cards.
            is_hidden = (provider == 'mistral')
            if is_hidden:
                print(f"[COUNCIL] Step {i+1}: {provider.upper()} (SHADOW EXECUTION) as {role.upper()}")
            else:
                print(f"[COUNCIL] Step {i+1}: {provider.upper()} as {role.upper()}")

            prompt = build_council_prompt(context, provider, role, i, total_steps)
            response_obj = {"success": False, "response": "Provider unknown"}

            # Primary Call — pass IDs for telemetry
            telemetry_kwargs = {"run_id": run_id, "session_id": session_id, "workflow": workflow, "user_id": user_id}
            try:
                if provider == 'openai': response_obj = call_openai_gpt4(prompt, role, images=images, **telemetry_kwargs)
                elif provider == 'anthropic': response_obj = call_anthropic_claude(prompt, role, images=images, **telemetry_kwargs)
                elif provider == 'google': response_obj = call_google_gemini(prompt, role, images=images, **telemetry_kwargs)
                elif provider == 'perplexity': response_obj = call_perplexity(prompt, role, **telemetry_kwargs)
                elif provider == 'mistral': response_obj = call_mistral_api(prompt, role, images=images, **telemetry_kwargs)
                elif provider == 'local': response_obj = call_local_llm(prompt, role, **telemetry_kwargs)
            except Exception as e:
                print(f"[COUNCIL] Primary ({provider}) Exception: {e}")
                response_obj = {"success": False}

            # 4. FALLBACKS (Omitted for brevity in REPLACEMENT but preserved in actual code logic)
            # [The fallback logic remains the same but should also pass telemetry_kwargs]
            # [I will keep the fallback logic and inject telemetry_kwargs below]
            
            # --- FALLBACK LEVEL 1: MISTRAL CLOUD ---
            if not response_obj.get('success', False) and provider != 'mistral':
                print(f"[COUNCIL] Primary ({provider}) Failed. Fallback to ANALYST (Mistral Cloud)...")
                try:
                    response_obj = call_mistral_api(prompt, role, **telemetry_kwargs)
                    if response_obj.get('success', False):
                        response_obj['response'] = f"[FALLBACK: ANALYST] {response_obj.get('response', '')}"
                        response_obj['model'] = f"Mistral (Fallback for {provider})"
                except Exception as e:
                    print(f"[COUNCIL] Mistral Fallback Exception: {e}")

            # --- FALLBACK LEVEL 2: LOCAL ORACLE ---
            if not response_obj.get('success', False) and provider != 'local':
                print(f"[COUNCIL] Secondary Failed. Fallback to ORACLE (Local)...")
                try:
                    response_obj = call_local_llm(prompt, "Oracle", **telemetry_kwargs)
                    if response_obj.get('success', False):
                        response_obj['response'] = f"[FALLBACK: ORACLE] {response_obj.get('response', '')}"
                        response_obj['model'] = "Local Oracle (Emergency)"
                except Exception as e:
                    print(f"[COUNCIL] Local Fallback Exception: {e}")
            
            response_text = "No response"
            usage = {}
            if isinstance(response_obj, dict):
                 response_text = response_obj.get('response', 'Error')
                 usage = response_obj.get('usage', {})
                 step_cost = usage.get('cost', 0.0)
                 total_run_cost += step_cost
                 total_latency_ms += usage.get('latency', 0)
                 print(f"[COST DEBUG] {provider}: cost={step_cost}, latency={usage.get('latency', 0)}, running_total={total_run_cost}")
                 if response_obj.get('success'):
                     models_used.append(provider)
            else:
                 response_text = str(response_obj)

            # Strip hallucinated Falcon-style tokens when Falcon is NOT active
            if not ghost_map:
                response_text = _strip_phantom_tokens(response_text)

            # --- KORUM NORMALIZATION LAYER (Section 5 Directive) ---
            response_text = normalization_layer(response_text)

            context.add_entry(provider, role, response_text, usage=usage)
            
            # --- SHADOW FILTER ---
            # If hidden (Mistral), do NOT add to the results dict returned to UI.
            if is_hidden:
                continue

            results[provider] = {
                "success": response_obj.get('success', False),
                "response": response_text,
                "model": response_obj.get('model', 'unknown') if isinstance(response_obj, dict) else 'unknown',
                "role": role.upper(),
                "usage": usage,
                "interrogations": [], # Placeholder for future features
                "verifications": [], # Placeholder for future features
                "error": response_obj.get('error') if not response_obj.get('success') else None
            }

            # --- LEDGER: model_reasoning ---
            if response_obj.get('success'):
                _ledger_write("model_reasoning", {
                    "model_name": response_obj.get('model', 'unknown') if isinstance(response_obj, dict) else 'unknown',
                    "model_role": role.upper(),
                    "latency_ms": usage.get('latency', 0),
                    "response_hash": _hl.sha256(response_text.encode()).hexdigest(),
                    "tokens_in": usage.get('input', 0),
                    "tokens_out": usage.get('output', 0),
                    "cost": usage.get('cost', 0.0),
                    "workflow": workflow,
                })

        # --- ACCOUNTABILITY & SCORES ---
        print(f"[COUNCIL] Audit Initiated: Verifying {len(results)} responses...")
        for provider in results:
            text = results[provider]['response']
            claims = identify_claims(text, provider, user_id=user_id)
            verified = verify_claims(claims, context.history)
            results[provider]['truth_meter'] = calculate_truth_score(verified)
            results[provider]['verified_claims'] = verified

        # --- DIVERGENCE ANALYSIS (Resilient) ---
        print(f"[COUNCIL] Resilience check: Identifying Divergence...")
        try:
            divergence = analyze_council_divergence(results, context, user_id=user_id)
        except Exception as div_error:
            print(f"[COUNCIL DIVERGENCE] Failed: {div_error}")
            divergence = _empty_divergence(f"Divergence analysis failed: {str(div_error)}")

    except Exception as e:
        print(f"[EXECUTION ERROR] {e}")
        # Always return a baseline divergence to prevent downstream crashes
        divergence = _empty_divergence(f"Execution failed: {str(e)}")
        # Don't return yet — allow the finalizer pass to attempt an emergency synthesis if results exist

    # 3b. FINAL ARBITER — Mathematical Consistency (Section 8a)
    print("[COUNCIL] Final Arbiter: Validating numerical consistency...")
    try:
        arbiter_report = final_arbiter(results, divergence)
        if arbiter_report['status'] != "CLEAN":
            print(f"[ARBITER] Status: {arbiter_report['status']} | "
                  f"Conflicts: {len(arbiter_report['conflicts'])} | "
                  f"Math violations: {len(arbiter_report['math_violations'])}")
        else:
            print(f"[ARBITER] Status: CLEAN — no numerical inconsistencies detected.")
    except Exception as arb_error:
        print(f"[ARBITER ERROR] Non-blocking: {arb_error}")
        arbiter_report = {"status": "ERROR", "narrative_directive": "", "conflicts": [],
                          "math_violations": [], "resolved_metrics": {}, "metric_inventory": {}}

    # 4. Hidden Finalizer Pass (Section 10 Directive)
    # If Mistral is present, it acts as a non-visible Finalizer to unify tone.
    # SKIP for RESEARCH workflow — clean report mode handles unification via synthesis + output filter.
    finalizer_response = None
    _skip_finalizer = workflow in ("RESEARCH", "TECH", "WAR_ROOM")
    if not _skip_finalizer and ('mistral' in active_models or 'mistral' in [m.split('-')[0].lower() for m in (execution_order if 'execution_order' in locals() else [])]):
        print("[COUNCIL] Hidden Finalizer Pass: Mistral unifying report tone...")
        try:
            # Build history_text for the hidden Finalizer pass (aggregates previous council findings)
            history_text = "\n\n".join([f"[{p.upper()}]: {r.get('response', '')}" for p, r in results.items() if isinstance(r, dict) and r.get('success')])

            # Mistral uses the workflow-specific goal to guide the final unification.
            # Local workflow variable check for safety
            wf_key = workflow if 'workflow' in locals() else "RESEARCH"
            dna = WORKFLOW_DNA.get(wf_key, WORKFLOW_DNA.get("RESEARCH", {}))
            outcome_goal = dna.get('goal', 'Professional Intelligence Brief')
            
            finalizer_prompt = f"""
            You are the KorumOS FINALIZER. Your mission is to unify the voice and structure of this intelligence product.
            
            CRITICAL STRENGTHS TO APPLY (Section 10 Directive):
            1. ADAPTIVE TONE (#4): Apply a consistent, authoritative, and mission-aligned voice. Remove persona-specific markers and unify the narrative into a single professional flow.
            2. DATA SYNTHESIS (#3): Identify critical data findings in the discussion. Ensure they are structured logically. If you see inconsistent formatting, normalize them into the [STRUCTURED_TABLE] model.
            
            RULES:
            - NO NEW CONCLUSIONS: Do not add metrics or findings not present in the discussion.
            - NO HEDGING: If the council made a decision, state it clearly.
            - FLOW: Ensure transitions between phases are logical and seamless.
            
            COUNCIL DISCUSSION:
            {history_text}
            
            Outcome Expected: {outcome_goal}
            """
            # Mistral runs as a background process, not a primary phase
            f_resp = call_mistral_api(finalizer_prompt, "Finalizer", user_id=user_id)
            if f_resp.get('success'):
                # Normalize the finalizer's work to catch any lingering MD
                finalizer_response = normalization_layer(f_resp.get('response', ''))
                print("[COUNCIL] Finalizer pass complete (Tone + Synthesis optimized).")
        except Exception as e:
            print(f"[COUNCIL FINALIZER ERROR] {e}")

    # 5. Synthesis (Final Guard)
    try:
        synthesis = synthesize_results(context, divergence_analysis=divergence, arbiter_report=arbiter_report, results=results, user_id=user_id)

        # --- FINAL RENDERING NORMALIZATION (Section 6 Directive) ---
        if isinstance(synthesis, dict):
            for s_id, s_content in synthesis.get('sections', {}).items():
                if isinstance(s_content, str):
                    synthesis['sections'][s_id] = normalization_layer(s_content)

            meta = synthesis.get('meta', {})
            if isinstance(meta, dict):
                if meta.get('summary') and isinstance(meta.get('summary'), str):
                    meta['summary'] = normalization_layer(meta['summary'])
                if meta.get('final_document') and isinstance(meta.get('final_document'), str):
                    meta['final_document'] = normalization_layer(meta['final_document'])

            if finalizer_response:
                if "meta" not in synthesis: synthesis["meta"] = {}
                synthesis["meta"]["finalizer_review"] = finalizer_response

    except Exception as syn_e:
        print(f"[SYNTHESIS CRASH] {syn_e}")
        synthesis = {
            "success": False,
            "error": str(syn_e),
            "meta": {"workflow": workflow if 'workflow' in locals() else "UNKNOWN"},
        }

    # --- LEDGER: council_synthesis ---
    _syn_meta = synthesis.get('meta', {}) if isinstance(synthesis, dict) else {}
    _ledger_write("council_synthesis", {
        "participating_models": list(results.keys()),
        "consensus_score": divergence.get('consensus_score') if isinstance(divergence, dict) else None,
        "composite_truth_score": _syn_meta.get('composite_truth_score'),
        "consensus_hash": _hl.sha256(
            (_syn_meta.get('summary', '') or '').encode()
        ).hexdigest(),
    })

    # --- CONTRIBUTION SCORING ---
    total_out_tokens = sum(r['usage'].get('output', 0) for r in results.values() if 'usage' in r)
    for p, r in results.items():
        if not r.get('success'):
            r['contribution_score'] = 0
            continue
        token_share = (r['usage'].get('output', 0) / total_out_tokens * 100) if total_out_tokens > 0 else 0
        truth_contribution = r.get('truth_meter', 50)
        verif_bonus = 20 if r.get('verified_claims') else 0
        r['contribution_score'] = int((token_share * 0.4) + (truth_contribution * 0.4) + verif_bonus)
        r['contribution_score'] = max(0, min(100, r['contribution_score']))

    # --- LEDGER: decision_outcome ---
    _composite = _syn_meta.get('composite_truth_score', 0) or 0
    _risk = "low" if _composite >= 70 else ("medium" if _composite >= 40 else "high")
    _ledger_write("decision_outcome", {
        "risk_score": _risk,
        "confidence": _composite,
        "supporting_models": list(set(models_used)),
        "total_cost": total_run_cost,
        "total_latency_ms": total_latency_ms,
        "workflow": workflow,
    })

    return {
        "consensus": f"COUNCIL ADJOURNED. Plan: {classification.get('outputType','report').upper()} generated via {len(results)} steps.",
        "results": results,
        "classification": classification,
        "synthesis": synthesis,
        "divergence": divergence,
        "arbiter": arbiter_report,
        "metrics": {
            "run_id": run_id,
            "session_id": session_id,
            "run_cost": total_run_cost,
            "latency_ms": total_latency_ms,
            "models_used": list(set(models_used)),
            "workflow": workflow
        }
    }

# --- ANALYTIC DIVERGENCE LAYER (The Comparator) ---
def analyze_council_divergence(results, context, user_id=None):
    """
    Compares council outputs to identify consensus, disagreement, and evidence gaps.
    Runs AFTER all council phases and accountability, BEFORE synthesis.
    Distinct from Red Team: this analyzes differences BETWEEN model conclusions.
    """
    # Build comparison payload from successful results
    model_outputs = {}
    for provider, result in results.items():
        if isinstance(result, dict) and result.get('success'):
            model_outputs[provider] = {
                "role": result.get("role", "unknown"),
                "response_excerpt": result.get("response", "")[:3000],
                "truth_meter": result.get("truth_meter", 50)
            }

    if len(model_outputs) < 2:
        return _empty_divergence("Insufficient council outputs for divergence analysis.")

    comparison_text = ""
    for provider, data in model_outputs.items():
        resp_text = data.get('response', '')
        excerpt = (resp_text[:1000] + '...') if len(resp_text) > 1000 else resp_text
        comparison_text += f"\n[{provider.upper()} — {data['role']}] (Truth: {data['truth_meter']}/100):\n{excerpt}\n"

    prompt = f"""
    You are an Analytic Divergence Engine. Your job is to compare multiple AI council outputs and identify where they AGREE, where they DISAGREE, and what evidence would resolve disagreements.

    This is NOT Red Team. You are not attacking the analysis — you are comparing conclusions across models.

    COUNCIL OUTPUTS:
    {comparison_text}

    ANALYZE THE FOLLOWING DIMENSIONS across all outputs:
    1. Core scenario assessment (what each model thinks is happening)
    2. Threat attribution (who/what is responsible)
    3. Confidence level (how certain each model is)
    4. Timeline/urgency (how immediate the threat/opportunity is)
    5. Priority action (what each model recommends doing first)
    6. Key assumptions (what each model takes for granted)

    SCORING RULES:
    - consensus_score: 0-100 (100 = all models perfectly aligned)
    - divergence_score: 0-100 (100 = maximum disagreement)
    - These should sum to approximately 100 (+/- 10 for nuance)
    - If 4/5+ models align on a core point → consensus_score >= 70
    - If confidence spread > 20 points across models → divergence_score += 15
    - If attribution differs materially → divergence_score += 20

    Return ONLY valid JSON (no markdown):
    {{
      "consensus_score": 75,
      "divergence_score": 25,
      "protocol_variance": false,
      "agreement_topics": [
        {{"topic": "Brief topic name", "detail": "What the models agree on", "confidence": "high|moderate|low", "providers": ["list of agreeing providers"]}}
      ],
      "contested_topics": [
        {{"topic": "Brief topic name", "positions": [{{"provider": "name", "position": "What this model concluded", "evidence": "Supporting evidence cited"}}], "severity": "critical|high|medium|low", "operational_impact": "How this disagreement affects decision-making"}}
      ],
      "confidence_gaps": [
        {{"description": "What is uncertain", "spread": "Range of confidence across models", "severity": "high|medium|low"}}
      ],
      "resolution_requirements": [
        {{"question": "What evidence would resolve this disagreement", "priority": "high|medium|low"}}
      ],
      "divergence_summary": "2-3 sentence summary of the overall divergence picture"
    }}

    IMPORTANT:
    - Set "protocol_variance" to true if divergence_score > 30
    - Be specific about WHICH providers disagree on WHAT
    - If models agree on everything, say so — don't manufacture disagreement
    """

    try:
        resp = call_openai_gpt4(prompt, "DivergenceAnalyst", model="gpt-4o-mini", user_id=user_id)
        if resp.get('success'):
            content = resp['response'].replace('```json', '').replace('```', '').strip()
            data = json.loads(content)
            # Ensure protocol_variance flag is set correctly
            data['protocol_variance'] = data.get('divergence_score', 0) > 30
            print(f"[DIVERGENCE] Consensus: {data.get('consensus_score', '?')}/100 | Divergence: {data.get('divergence_score', '?')}/100 | Variance: {data.get('protocol_variance', False)}")
            return data
    except Exception as e:
        print(f"[DIVERGENCE ERROR] {e}")

    return _empty_divergence("Divergence analysis failed — proceeding with synthesis.")


def _empty_divergence(reason=""):
    """Returns a safe empty divergence structure."""
    return {
        "consensus_score": 50,
        "divergence_score": 0,
        "protocol_variance": False,
        "agreement_topics": [],
        "contested_topics": [],
        "confidence_gaps": [],
        "resolution_requirements": [],
        "divergence_summary": reason or "Divergence analysis not available."
    }




# ---------------------------------------------------------------------------
# MIMIR STRUCTURED AUDIT PROMPT
# ---------------------------------------------------------------------------
# Injected into the council prompt when a Ghost Map + residual report is
# present. Forces models to operate as forensic auditors instead of
# generic commentators — they MUST report from the token inventory,
# not hallucinate about a document they "don't have access to."
# ---------------------------------------------------------------------------

def _build_mimir_block(ghost_map: dict, residual_report: dict) -> str:
    """
    Build the MIMIR structured audit block to prepend to council prompts
    when analyzing a Falcon-redacted document.

    Args:
        ghost_map: output of build_ghost_map_summary() — safe token inventory
        residual_report: output of detect_residual_pii() — what Falcon missed

    Returns:
        A formatted string block for injection into the council prompt.
    """
    if not ghost_map:
        return ""

    token_inventory = ghost_map.get("token_inventory", [])
    by_type = ghost_map.get("by_type", {})
    total = ghost_map.get("total_redacted", 0)
    high_risk = ghost_map.get("high_risk_types", [])
    falcon_level = ghost_map.get("falcon_level", "STANDARD")

    residual_count = residual_report.get("residual_count", 0)
    residuals = residual_report.get("residuals", [])
    audit_note = residual_report.get("audit_note", "")

    # Build token inventory table (capped at 40 entries to stay within context)
    inv_lines = []
    for i, entry in enumerate(token_inventory[:40]):
        inv_lines.append(
            f"  {entry['token']:<20} | {entry['entity_type']:<14} | "
            f"{entry.get('raw_category', entry['entity_type']):<16} | "
            f"pass={entry.get('source_pass','?')}"
        )
    if len(token_inventory) > 40:
        inv_lines.append(f"  ... and {len(token_inventory) - 40} more tokens (truncated for context)")

    # Build type summary
    type_lines = []
    for etype, tokens in sorted(by_type.items()):
        type_lines.append(f"  {etype:<16}: {len(tokens):>3} token(s) — {', '.join(tokens[:6])}{'...' if len(tokens) > 6 else ''}")

    # Build residual section
    if residual_count == 0:
        residual_block = "  PII_DIFF_CLEAN: No residual PII detected after primary Falcon pass."
    else:
        res_lines = []
        for r in residuals[:10]:
            _conf = r["confidence"].upper()
            _cat = r["category"]
            _frag = r["text_fragment"][:60]
            _off = r["char_offset"]
            res_lines.append(f'  [{_conf}] {_cat}: "{_frag}" @ offset {_off}')
        if len(residuals) > 10:
            res_lines.append(f"  ... and {len(residuals) - 10} more residuals")
        residual_block = "\n".join(res_lines)

    block = f"""
================================================================================
MIMIR PROTOCOL — FALCON AUDIT CONTEXT
You are operating as a FORENSIC PII AUDITOR, not a general commentator.
The document you are analyzing has been pre-processed by the Falcon redaction
engine at level: {falcon_level}.

You MUST ground ALL of your analysis in the Ghost Map and PII Diff data below.
DO NOT say "I would need access to the document" — the token inventory IS the
document structure. Reason from what is here.

── GHOST MAP: TOKEN INVENTORY ({total} entities redacted) ─────────────────────
  TOKEN                | ENTITY_TYPE    | RAW_CATEGORY     | SOURCE_PASS
  ─────────────────────┼────────────────┼──────────────────┼─────────────────
{chr(10).join(inv_lines)}

── ENTITY TYPE SUMMARY ────────────────────────────────────────────────────────
{chr(10).join(type_lines) if type_lines else "  No entities detected."}

── HIGH-RISK CATEGORIES ────────────────────────────────────────────────────────
  {', '.join(high_risk) if high_risk else 'None identified.'}

── PII DIFF — RESIDUAL DETECTION ({residual_count} items missed) ───────────────
{residual_block}

── AUDIT NOTE ──────────────────────────────────────────────────────────────────
  {audit_note}

── YOUR MANDATORY OUTPUT FORMAT ────────────────────────────────────────────────
You MUST structure your response to include ALL of the following sections:

1. GHOST MAP REPORT
   List every token from the inventory above. For each, state:
   - Token ID (e.g. [PERSON_01])
   - Entity type and inferred role in the document (e.g. "natural_person_name, General Counsel")
   - Clause or section where it appears (if inferrable from context)

2. CLAUSE-BY-CLAUSE PII AUDIT
   For each major document section visible in the redacted text:
   - State whether residual (missed) PII is present: YES / NO
   - If YES: cite the specific residual item and its category

3. RE-IDENTIFICATION FEASIBILITY OPINION
   Rate: LOW / MEDIUM / HIGH
   Provide specific reasoning, e.g.:
   "HIGH — attacker can link [ORG_01] + Project Sponsor role to identify [PERSON_01]
   via public SEC filings. [PERSON_02] phone number fragment (SSN_PARTIAL) narrows
   identity to 3 individuals matching other document context."

4. RECOMMENDED REMEDIATION
   Specific Falcon configuration changes or additional redaction passes needed.
   Reference actual token IDs and categories from the Ghost Map.

If the ghost map is empty or the document has no redactions, state:
"GHOST MAP EMPTY — document may contain no PII, or Falcon was not applied.
Recommend re-running at STANDARD level before council analysis."
================================================================================
"""
    return block


def build_council_prompt(context, ai_name, persona, position, total_steps):
    # Determine the task objective
    core_objective = context.query
    intent = context.classification.get('intent', 'analysis')
    output_type = context.classification.get('outputType', 'report')

    # Retrieve Workflow DNA
    dna = WORKFLOW_DNA.get(context.workflow, WORKFLOW_DNA["RESEARCH"])
    if "alias" in dna:
        dna = WORKFLOW_DNA.get(dna["alias"], dna)

    # --- DECISION ENFORCEMENT (Section 8 Directive) ---
    decision_enforcement = """
    ## DECISION ENFORCEMENT (Section 8 Directive)
    - DO NOT use hedging language: "might", "could", "possibly", "it depends".
    - BE DECISIVE. No 50/50 answers. No generic output.
    - Every finding or recommendation MUST include:
        1. DECISION: A clear yes/no/proceed/halt choice.
        2. ACTION: The immediate next step to take.
        3. CONFIDENCE: High/Medium/Low with rationale.
    """

    # --- DATA INTEGRITY (Section 8a Directive) ---
    data_integrity = """
    ## DATA INTEGRITY (Section 8a Directive)
    - Do NOT introduce external standards, regulations, or compliance references unless explicitly provided in the dataset or prompt.
    - Do NOT fabricate legal, tax, or accounting requirements.
    - All conclusions must be directly supported by the provided data.
    - Do NOT cite regulations, frameworks, or laws that are not present in the input.
    - If no data supports a claim, state that explicitly — do NOT invent supporting evidence.
    """

    # Select workflow-specific directives if available, else generic
    active_phase_directives = WORKFLOW_PHASE_OVERRIDES.get(context.workflow, PHASE_DIRECTIVES)


    if ai_name.lower() == 'perplexity':
        # Perplexity works best with direct research questions, not persona roleplay
        prompt = f"""
        RESEARCH OBJECTIVE: "{core_objective}"

        CONTEXT: You are contributing to a strategic intelligence report (Phase {position + 1} of {total_steps}).
        GOAL: {dna['goal']}
        INTENT: {intent.upper()} / {output_type.upper()}
        """

        if position > 0:
            prompt += "\nRECENT DATA / PREVIOUS ANALYSIS:\n"
            last_entry = context.history[-1]
            prompt += f"Summary of previous findings: {last_entry['response'][:1000]}...\n"

        # Give Perplexity the phase-specific focus even in research mode
        phase = active_phase_directives.get(position, active_phase_directives.get(min(position, 4)))
        prompt += f"\nYOUR SPECIFIC FOCUS: {phase['title']}\n{phase['instruction']}"
        prompt += decision_enforcement # Section 8
        prompt += data_integrity # Section 8a
        if context.ghost_map or context.residual_report:
            prompt += "\n" + _build_mimir_block(context.ghost_map, context.residual_report)
        prompt += "\nProvide comprehensive, well-sourced research and data. Focus on facts, metrics, and technical details."
        prompt += "\nNEVER invent bracket-notation placeholders like [CURRENCY_AMOUNT_XXXX] or [ENTITY_TYPE_HASH]. Use real values from the source material. If unknown, say so."
        return prompt

    # --- Determine which phase directive applies ---
    # Map position to phase (if more than 5 steps, later steps get the closest directive)
    if position == total_steps - 1 and total_steps >= 3:
        # Final step is ALWAYS the integrator (phase 4) regardless of count
        phase = active_phase_directives[4]
    elif total_steps <= 5:
        phase = active_phase_directives.get(position, active_phase_directives[min(position, 4)])
    else:
        # Scale phases across available steps
        phase_index = int(position / total_steps * 5)
        phase = active_phase_directives.get(min(phase_index, 4))

    # Inject MIMIR block if Falcon ghost map is present
    mimir_block = ""
    if context.ghost_map or context.residual_report:
        mimir_block = _build_mimir_block(context.ghost_map, context.residual_report)

    # Document-assembly workflows: final phase MUST compile everything, not just add unique content
    DOCUMENT_ASSEMBLY_WORKFLOWS = {"EOM_STATEMENT", "FINANCE", "AUDIT", "CODE_AUDIT", "LEGAL", "PORTFOLIO_BUILDER"}
    is_final_phase = (position == total_steps - 1)
    is_assembly_workflow = context.workflow in DOCUMENT_ASSEMBLY_WORKFLOWS

    if is_final_phase and is_assembly_workflow:
        assembly_rule = (
            "## CRITICAL RULE: FULL DOCUMENT ASSEMBLY\n"
            "    You are the FINAL AUTHOR. Your job is to compile ALL prior phase outputs into ONE complete, "
            "ready-to-present document. Use every figure, table, and finding from the prior phases — "
            "structured exactly as your phase mission defines. This is NOT a summary. It is the final deliverable. "
            "If it is not complete enough to hand to a client or executive, you have FAILED."
        )
    else:
        assembly_rule = (
            "## CRITICAL RULE: ZERO REPETITION\n"
            "    The prior phases are provided for CONTEXT ONLY. DO NOT restate, summarize, or echo their findings.\n"
            "    Your ONLY job is to add what is MISSING \u2014 the unique contribution defined by your phase mission above.\n"
            "    If your output overlaps with prior phases, you have FAILED your mission."
        )

    # Standard Professional Template with DNA Overlay
    prompt = f"""
    ## PROFESSIONAL INTELLIGENCE BRIEF
    MISSION TYPE: {context.workflow}
    PRIMARY OBJECTIVE: "{core_objective}"

    CURRENT FOCUS: {persona.upper()} (Assignee: {ai_name.upper()})
    WORKFLOW STAGE: Phase {position + 1} of {total_steps} — {phase['title']}

    --- WORKFLOW DNA ---
    POSTURE: {dna['posture']}
    GOAL: {dna['goal']}
    TONE: {dna['tone']}
    RISK BIAS: {dna['risk_bias']}
    TIME HORIZON: {dna['time_horizon']}
    REPORT SECTIONS: {', '.join(dna['output_structure'])}
    --------------------

    ## YOUR PHASE MISSION
    {phase['instruction']}

    {assembly_rule}

    {mimir_block}
    ## INTERNAL STRUCTURING (FOR SYSTEM PARSING)
    You MUST use the following tags to mark high-value intelligence:
    - [DECISION_CANDIDATE] ...recommendation text... [/DECISION_CANDIDATE]
    - [RISK_VECTOR] ...risk description... [/RISK_VECTOR]
    - [METRIC_ANCHOR] ...metric value... [/METRIC_ANCHOR]
    - [TRUTH_BOMB] ...critically verified fact... [/TRUTH_BOMB]

    Do not let these tags disrupt your narrative flow; they are for the backend extractor.

    ## STRICT OUTPUT RULE
    ONLY use the four system tags listed above ([DECISION_CANDIDATE], [RISK_VECTOR], [METRIC_ANCHOR], [TRUTH_BOMB]).
    NEVER invent bracket-notation placeholders such as [CURRENCY_AMOUNT_XXXX], [ENTITY_TYPE_HASH], or any
    similar token. Always use real values, names, and figures from the source material. If a value is unknown,
    say "unknown" or "not provided" — do NOT fabricate redaction-style placeholder tokens.
    
    {decision_enforcement}

    {data_integrity}
    """

    # Inject prior session context for follow-up queries
    if context.previous_context:
        prompt += "\n## PRIOR SESSION CONTEXT (Follow-up Query)\n"
        prompt += "The council previously analyzed a related query. Build on these conclusions — do NOT repeat them.\n"
        for entry in context.previous_context[-2:]:
            prompt += f"\n**Previous Query:** \"{entry.get('query', 'N/A')}\"\n"
            prompt += f"**Prior Consensus ({entry.get('consensus_score', 'N/A')}/100):** {entry.get('summary', 'N/A')}\n"
            contested = entry.get('contested_topics', [])
            if contested:
                prompt += f"**Unresolved Disputes:** {', '.join(contested)}\n"
            div_summary = entry.get('divergence_summary', '')
            if div_summary and div_summary != 'Divergence analysis not available.':
                prompt += f"**Divergence:** {div_summary}\n"
            # Full Loop Integrity: include interrogation/verification history
            interrog_history = entry.get('interrogation_history', '')
            if interrog_history:
                prompt += f"\n**Prior Interrogation/Verification Events:**\n{interrog_history}\n"
                prompt += "IMPORTANT: These events represent challenges to the prior council output. "
                prompt += "Your response MUST account for any concessions or invalidations noted above. "
                prompt += "Do NOT repeat claims that were previously invalidated.\n"
        prompt += "\n--------------------\n"

    if position == 0:
        prompt += f"\n## ASSIGNMENT:\nBuild the neutral intake baseline for this {context.workflow} mission. Facts only. No opinions. No strategy."
        if context.previous_context:
            prompt += "\nNote: A prior session's conclusions are provided above. Acknowledge them briefly, then focus on what is NEW in this follow-up query."
    else:
        # Final assembly phase gets full prior phase output; others get 2000-char snippets
        max_chars = 6000 if (is_final_phase and is_assembly_workflow) else 2000
        header = "\n## PRIOR PHASE OUTPUT (COMPILE INTO FINAL DOCUMENT):\n" if (is_final_phase and is_assembly_workflow) else "\n## PRIOR PHASE CONTEXT (for reference — DO NOT REPEAT):\n"
        prompt += header
        for entry in context.history:
            snippet = (entry['response'][:max_chars] + '...') if len(entry['response']) > max_chars else entry['response']
            prompt += f"\n-- PHASE [{entry['persona'].upper()}] ({entry['ai'].upper()}):\n{snippet}\n"

    return prompt

# ===================================================================
# CONFIDENCE CALIBRATION — Score-Band Language Enforcement
# ===================================================================

CONFIDENCE_BANDS = [
    (80, 100, "high", "Strong evidence supports conclusions. Language may use 'confirms', 'demonstrates', 'establishes'."),
    (60,  79, "moderate-to-high", "Moderate confidence. Use 'suggests', 'indicates', 'likely'. Avoid absolute certainty."),
    (40,  59, "moderate", "Weak support. Use 'preliminary data suggests', 'may indicate', 'possible'. Emphasize caveats."),
    ( 0,  39, "low", "Insufficient data. Use 'unconfirmed', 'speculative', 'insufficient evidence to determine'. Recommend further investigation."),
]

def _build_confidence_directive(results):
    """
    Compute average truth score from results and return the appropriate
    confidence band directive for injection into the synthesis prompt.
    """
    if not results:
        return "Default to 'moderate-to-high' confidence. Only use 'high' if 3+ council members independently confirmed facts with specific evidence. This is an intelligence product — overconfidence is a liability."

    scores = [r.get('truth_meter', 50) for r in results.values() if isinstance(r, dict) and r.get('success')]
    if not scores:
        avg = 50
    else:
        avg = sum(scores) / len(scores)

    band_label = "moderate-to-high"
    band_guidance = ""
    for low, high, label, guidance in CONFIDENCE_BANDS:
        if low <= avg <= high:
            band_label = label
            band_guidance = guidance
            break

    directive = (
        f"The council's aggregate truth score is {avg:.0f}/100.\n"
        f"    This places the output in the '{band_label.upper()}' confidence band.\n"
        f"    Required tone: {band_guidance}\n"
        f"    You MUST set overall_confidence to '{band_label}'.\n"
        f"    Your language throughout ALL sections must match this band.\n"
        f"    VIOLATION: Using 'confirms' or 'demonstrates' when score is below 80 is overconfidence.\n"
        f"    VIOLATION: Using 'speculative' or 'insufficient' when score is above 79 is underconfidence.\n"
        f"    Match your language to the evidence. This is decision intelligence — calibration is everything."
    )
    return directive


def _is_decision_packet_shape(data):
    if not isinstance(data, dict):
        return False
    packet_keys = {
        "decision_headline",
        "executive_summary",
        "go_no_go_call",
        "confidence",
        "verified_claims",
        "risk_vectors",
        "assumptions",
        "unknowns",
        "immediate_actions",
        "alternatives_rejected",
        "evidence_trace",
        "export_metadata",
    }
    present = sum(1 for key in packet_keys if key in data)
    return present >= 8 and "sections" not in data


def _packet_confidence_status(score):
    if score >= 80:
        return "Recommended"
    if score >= 70:
        return "Conditional"
    return "Fail"


def _verified_fact_count(verified_claims):
    count = 0
    for claim in verified_claims or []:
        if not isinstance(claim, dict):
            continue
        status = str(claim.get("status") or "").strip().upper()
        if status in ("VERIFIED", "ACCURATE"):
            count += 1
    return count


def _critical_data_missing(verified_claims, evidence_trace, unknowns):
    verified_count = _verified_fact_count(verified_claims)
    evidence_count = len([item for item in (evidence_trace or []) if str(item).strip()])
    if verified_count == 0 or evidence_count == 0:
        return True

    unknown_text = " ".join(str(item).strip().lower() for item in (unknowns or []) if str(item).strip())
    critical_markers = (
        "baseline", "probability", "projection", "forecast", "cost", "budget",
        "timeline", "frequency", "churn", "latency", "outage", "failure threshold",
    )
    return any(marker in unknown_text for marker in critical_markers)


def _calibrate_packet_confidence(packet, score, basis, red_team_findings=None):
    """
    Confidence Governor — enforces dual-confidence model with hard caps.

    Returns: (score, status, basis, fact_confidence, decision_confidence)

    fact_confidence   = how certain we are about the root cause (0.0–1.0)
    decision_confidence = how confident we are in the ACTION (0.0–1.0)

    Governor rules cap these independently. Both cannot be HIGH without strong evidence.
    """
    verified_claims = packet.get("verified_claims") or []
    evidence_trace = packet.get("evidence_trace") or []
    unknowns = [str(item).strip() for item in (packet.get("unknowns") or []) if str(item).strip()]
    assumptions = [str(item).strip() for item in (packet.get("assumptions") or []) if str(item).strip()]
    verified_count = _verified_fact_count(verified_claims)
    unknown_count = len(unknowns)
    assumption_count = len(assumptions)
    quantified = _has_quantified_support(verified_claims, evidence_trace)
    critical_missing = _critical_data_missing(verified_claims, evidence_trace, unknowns)
    notes = []

    # --- FACT CONFIDENCE GOVERNOR ---
    fact_confidence = 0.8  # Start optimistic, apply caps

    # No telemetry/metrics in evidence → max 0.5
    if not quantified:
        fact_confidence = min(fact_confidence, 0.5)
        notes.append("no quantified telemetry in evidence")

    # Root cause not directly observed → max 0.3
    if critical_missing:
        fact_confidence = min(fact_confidence, 0.3)
        notes.append("root cause not directly observed")

    # No verified facts at all → max 0.2
    if verified_count == 0:
        fact_confidence = min(fact_confidence, 0.2)
        notes.append("zero verified facts")

    # --- DECISION CONFIDENCE GOVERNOR ---
    decision_confidence = 0.8  # Start optimistic, apply deductions

    # Per major assumption → subtract 0.05–0.10
    if assumption_count > 0:
        assumption_penalty = min(assumption_count * 0.07, 0.35)
        decision_confidence -= assumption_penalty
        notes.append(f"{assumption_count} assumptions penalize confidence")

    # Unknowns > 3 → subtract 0.10
    if unknown_count > 3:
        decision_confidence -= 0.10
        notes.append("more than 3 unknowns")

    # Unknowns exceed verified facts → hard cap
    if unknown_count > verified_count:
        decision_confidence = min(decision_confidence, 0.6)
        notes.append("unknowns exceed verified facts")

    # Red Team findings apply additional pressure
    if red_team_findings and isinstance(red_team_findings, dict):
        rt_status = str(red_team_findings.get("red_team_status", "")).upper()
        # Confidence attack → subtract 0.10 (direct challenge to decision certainty)
        conf_attack = str(red_team_findings.get("confidence_attack", "")).strip()
        if conf_attack and len(conf_attack) >= 20:
            decision_confidence -= 0.10
            notes.append("Red Team confidence attack valid")
        # Viable alternative strategy → subtract 0.15
        alt = str(red_team_findings.get("alternative_strategy", "")).strip()
        if alt and len(alt) >= 20:
            decision_confidence -= 0.15
            notes.append("Red Team identified viable alternative")
        # Missing evidence items reduce fact confidence
        missing = red_team_findings.get("missing_evidence") or []
        if len(missing) >= 2:
            fact_confidence = min(fact_confidence, fact_confidence - 0.10)
            notes.append("Red Team identified missing evidence")
        # Unsupported claims reduce decision confidence
        unsupported = red_team_findings.get("unsupported_claims") or []
        if len(unsupported) >= 2:
            decision_confidence -= 0.10
            notes.append("Red Team identified unsupported claims")
        # FAIL status → hard caps
        if "FAIL" in rt_status:
            decision_confidence = min(decision_confidence, 0.5)
            notes.append("Red Team status is FAIL")

    # Clamp to [0.0, 1.0]
    fact_confidence = max(0.0, min(1.0, round(fact_confidence, 2)))
    decision_confidence = max(0.0, min(1.0, round(decision_confidence, 2)))

    # Never allow both to be HIGH without strong evidence
    if fact_confidence >= 0.8 and decision_confidence >= 0.8:
        if not quantified or critical_missing or unknown_count > verified_count:
            decision_confidence = min(decision_confidence, 0.7)
            notes.append("dual-high blocked without strong evidence")

    # --- ROOT CAUSE GATE ---
    # No root cause established → NO action/upgrade language. Decision must be CONDITIONAL.
    go_no_go = packet.get("go_no_go_call") or {}
    decision_text = str(go_no_go.get("decision") or "").strip().upper()
    if fact_confidence <= 0.3:
        if decision_text == "GO":
            go_no_go["decision"] = "CONDITIONAL GO"
            notes.append("root cause gate: no root cause established, action language blocked")
        if go_no_go.get("rationale"):
            go_no_go["rationale"] = f"[ROOT CAUSE UNESTABLISHED] {go_no_go['rationale']}"

    # --- ASSUMPTION FIREWALL ---
    # If decision depends on unverified assumptions → force CONDITIONAL
    decision_text = str(go_no_go.get("decision") or "").strip().upper()  # Re-read after root cause gate
    if decision_text == "GO" and assumption_count > 0 and verified_count < assumption_count:
        # More assumptions than verified facts → cannot be GO
        go_no_go["decision"] = "CONDITIONAL GO"
        notes.append("assumption firewall: GO blocked, assumptions exceed verified facts")

    # --- MAP TO LEGACY SCORE ---
    # Legacy score (0-100) = weighted average of both confidences
    combined = (fact_confidence * 0.4 + decision_confidence * 0.6) * 100
    score = min(score, int(round(combined)))

    if notes:
        note_text = "Confidence Governor: " + "; ".join(notes) + "."
        basis = f"{basis} {note_text}".strip() if basis else note_text

    return score, _packet_confidence_status(score), basis, fact_confidence, decision_confidence


def _has_quantified_support(verified_claims, evidence_trace):
    texts = []
    for claim in verified_claims or []:
        if isinstance(claim, dict):
            texts.extend([
                str(claim.get("claim") or ""),
                str(claim.get("source_ref") or ""),
            ])
        else:
            texts.append(str(claim))
    for item in evidence_trace or []:
        if isinstance(item, dict):
            texts.extend([
                str(item.get("point") or ""),
                str(item.get("detail") or ""),
                str(item.get("source") or ""),
            ])
        else:
            texts.append(str(item))
    blob = " ".join(t for t in texts if t).lower()
    return bool(re.search(r"\b\d+(?:\.\d+)?%?\b", blob))


def _requires_diagnostic_first(packet):
    verified_claims = packet.get("verified_claims") or []
    evidence_trace = packet.get("evidence_trace") or []
    unknowns = [str(item).strip() for item in (packet.get("unknowns") or []) if str(item).strip()]
    verified_count = _verified_fact_count(verified_claims)
    evidence_count = len([item for item in evidence_trace if str(item).strip()])
    unknown_count = len(unknowns)
    quantified_support = _has_quantified_support(verified_claims, evidence_trace)
    critical_missing = _critical_data_missing(verified_claims, evidence_trace, unknowns)

    reasons = []
    if verified_count < 2:
        reasons.append("fewer than two verified facts")
    if evidence_count < 2:
        reasons.append("fewer than two supporting evidence points")
    if unknown_count > verified_count:
        reasons.append("unknowns outweigh verified facts")
    if critical_missing:
        reasons.append("critical baseline data is missing")
    if not quantified_support and critical_missing:
        reasons.append("no quantified baseline support is available")

    gate = (
        (verified_count < 2 and evidence_count < 2)
        or (critical_missing and unknown_count >= verified_count)
        or (critical_missing and not quantified_support)
    )
    return gate, reasons


def _packet_timeline_bucket(timeline):
    tl = str(timeline or "").strip().lower()
    if any(term in tl for term in ("immediate", "now", "urgent", "day 0")):
        return "Immediate"
    if any(term in tl for term in ("mid", "later", "long", "90", "quarter", "phase 3")):
        return "Mid-Term"
    if any(term in tl for term in ("near", "soon", "next", "30", "60", "phase 2")):
        return "Near-Term"
    return "Immediate"


def _packet_bullet_lines(items):
    return "\n".join(f"- {item}" for item in items if str(item).strip())


def _group_actions_from_packet(actions):
    grouped = {"Immediate": [], "Near-Term": [], "Mid-Term": []}
    for action in actions or []:
        if isinstance(action, dict):
            line = str(action.get("action") or "").strip()
            owner = str(action.get("owner") or "").strip()
            if owner:
                line = f"{line} ({owner})"
            bucket = _packet_timeline_bucket(action.get("timeline"))
        else:
            line = str(action).strip()
            bucket = "Immediate"
        if line:
            grouped.setdefault(bucket, []).append(line)

    parts = []
    for label in ("Immediate", "Near-Term", "Mid-Term"):
        if grouped.get(label):
            parts.append(f"**{label}**\n" + _packet_bullet_lines(grouped[label]))
    return "\n\n".join(parts)


def adapt_decision_packet_to_legacy_shape(packet, workflow="RESEARCH"):
    """
    Convert a strict Decision Packet into the legacy synthesis shape still
    required by exporters and clean-report validation.
    """
    if not isinstance(packet, dict):
        return packet

    export_meta = packet.get("export_metadata") or {}
    confidence = packet.get("confidence") or {}
    go_no_go = packet.get("go_no_go_call") or {}
    verified_claims = packet.get("verified_claims") or []
    evidence_trace = packet.get("evidence_trace") or []
    risk_vectors = packet.get("risk_vectors") or []
    assumptions = [str(item).strip() for item in (packet.get("assumptions") or []) if str(item).strip()]
    unknowns = [str(item).strip() for item in (packet.get("unknowns") or []) if str(item).strip()]
    alternatives = packet.get("alternatives_rejected") or []
    actions = packet.get("immediate_actions") or []

    raw_score = confidence.get("score", export_meta.get("composite_truth_score", 82))
    try:
        score = int(round(float(raw_score)))
    except (TypeError, ValueError):
        score = 82
    basis = str(confidence.get("basis") or "").strip()
    red_team_findings = packet.get("_red_team_findings")
    score, status, basis, fact_confidence, decision_confidence = _calibrate_packet_confidence(
        packet, score, basis, red_team_findings=red_team_findings
    )
    diagnostic_first, diagnostic_reasons = _requires_diagnostic_first(packet)

    decision_text = str(go_no_go.get("decision") or "CONDITIONAL GO").strip().upper()
    rationale_text = str(go_no_go.get("rationale") or "").strip()
    if diagnostic_first and decision_text == "GO":
        decision_text = "CONDITIONAL GO"
    if diagnostic_first:
        diagnostic_note = "INSUFFICIENT EVIDENCE for root-cause attribution. Collect baseline diagnostics before committing to irreversible remediation."
        rationale_text = f"{diagnostic_note} {rationale_text}".strip() if rationale_text else diagnostic_note
        score = min(score, 65)
        if diagnostic_reasons:
            basis_note = "Diagnostic-first gating applied because " + ", ".join(dict.fromkeys(diagnostic_reasons)) + "."
            basis = f"{basis} {basis_note}".strip() if basis else basis_note
    if diagnostic_first and decision_text == "CONDITIONAL GO":
        score = 70
    status = _packet_confidence_status(score)

    key_signals = []
    for claim in verified_claims[:4]:
        if isinstance(claim, dict):
            claim_text = str(claim.get("claim") or "").strip()
        else:
            claim_text = str(claim).strip()
        if claim_text:
            key_signals.append(claim_text)

    if len(key_signals) < 4:
        for trace in evidence_trace:
            if isinstance(trace, dict):
                point = str(trace.get("point") or "").strip()
            else:
                point = str(trace).strip()
            if point and point not in key_signals:
                key_signals.append(point)
            if len(key_signals) >= 4:
                break

    high_impact = []
    secondary = []
    structured_risks = []
    for risk in risk_vectors:
        if isinstance(risk, dict):
            title = str(risk.get("title") or "Risk").strip()
            description = str(risk.get("description") or "").strip()
            mitigation = str(risk.get("mitigation") or "").strip()
            severity = str(risk.get("severity") or "MEDIUM").strip().upper()
            line = f"{title}: {description}" if description else title
            if mitigation:
                line += f" Mitigation: {mitigation}"
            if severity in ("CRITICAL", "HIGH"):
                high_impact.append(line)
            else:
                secondary.append(line)
            structured_risks.append({
                "risk": title,
                "severity": severity.lower(),
                "mitigation": mitigation or "None specified",
            })
        else:
            text = str(risk).strip()
            if text:
                secondary.append(text)

    critical_challenges_parts = []
    if high_impact:
        critical_challenges_parts.append("**HIGH IMPACT CHALLENGES**\n" + _packet_bullet_lines(high_impact))
    if secondary:
        critical_challenges_parts.append("**SECONDARY CONSIDERATIONS**\n" + _packet_bullet_lines(secondary))
    critical_challenges_parts.append(
        "**SYNTHESIS**: The uncertainties do not outweigh the risk of inaction; proceed while monitoring the key assumptions and unknowns."
    )

    tradeoff_lines = []
    for alt in alternatives:
        if isinstance(alt, dict):
            option = str(alt.get("option") or "").strip()
            reason = str(alt.get("reason_rejected") or "").strip()
            if option or reason:
                tradeoff_lines.append(f"- {option}: {reason}".strip(": "))
        else:
            text = str(alt).strip()
            if text:
                tradeoff_lines.append(f"- {text}")

    execution_lines = []
    if assumptions:
        execution_lines.append("**ASSUMPTIONS**\n" + _packet_bullet_lines(assumptions))
    if unknowns:
        execution_lines.append("**UNKNOWNS**\n" + _packet_bullet_lines(unknowns))

    confidence_lines = [
        f"Score: {score} / 100",
        f"Status: {status}",
    ]
    if basis:
        confidence_lines.extend(["", "Key Drivers:", f"- {basis}"])
    if assumptions:
        confidence_lines.extend(["", "Assumptions:"])
        confidence_lines.extend(f"- {item}" for item in assumptions[:5])
    if unknowns:
        confidence_lines.extend(["", "Limitations:"])
        confidence_lines.extend(f"- {item}" for item in unknowns[:5])

    rationale_bullets = []
    if rationale_text:
        rationale_bullets.append(rationale_text)
    rationale_bullets.extend(key_signals[:4])
    decision_lines = [f"Decision: {decision_text}"]
    if rationale_bullets:
        decision_lines.extend(f"- {item}" for item in rationale_bullets[:5])

    final_assessment = str(
        packet.get("decision_headline")
        or rationale_text
        or packet.get("executive_summary")
        or ""
    ).strip()
    if not final_assessment:
        final_assessment = "Proceed on the strongest defensible path while monitoring the exposed risks and unresolved assumptions."
    if diagnostic_first:
        final_assessment = "Do not commit to irreversible remediation until baseline diagnostics confirm the primary constraint."

    summary_text = str(packet.get("executive_summary") or "").strip() or final_assessment
    if diagnostic_first:
        summary_text = (
            "Current evidence is insufficient to attribute the problem to a single root cause with confidence. "
            "Leadership should hold major remediation to a diagnostic-first path until baseline telemetry and operating constraints are verified. "
            "The recommendation remains conditional because unknowns outweigh the verified evidence."
        )

    legacy = {
        "meta": {
            "title": str(packet.get("decision_headline") or "Decision Packet").strip() or "Decision Packet",
            "generated_at": export_meta.get("generated_at") or datetime.now().isoformat(),
            "summary": summary_text,
            "workflow": export_meta.get("workflow") or workflow,
            "composite_truth_score": score,
            "truth_score": score,
            "fact_confidence": fact_confidence,
            "decision_confidence": decision_confidence,
        },
        "sections": {
            "executive_summary": summary_text,
            "key_signals": _packet_bullet_lines(key_signals),
            "critical_challenges": "\n\n".join(part for part in critical_challenges_parts if part),
            "tradeoffs": "\n".join(tradeoff_lines),
            "decision": "\n".join(decision_lines),
            "action_priorities": _group_actions_from_packet(actions),
            "execution_considerations": "\n\n".join(execution_lines),
            "confidence_assessment": "\n".join(confidence_lines),
            "final_assessment": final_assessment,
        },
        "structured_data": {
            "key_metrics": [
                {"metric": "Decision Score", "value": f"{score} / 100", "context": status},
                {"metric": "Fact Confidence", "value": f"{fact_confidence}", "context": "Root cause certainty"},
                {"metric": "Decision Confidence", "value": f"{decision_confidence}", "context": "Action certainty"},
            ],
            "action_items": [
                {
                    "task": str(action.get("action") or "").strip(),
                    "priority": "high" if _packet_timeline_bucket(action.get("timeline")) == "Immediate" else "medium",
                    "timeline": _packet_timeline_bucket(action.get("timeline")),
                }
                for action in actions
                if isinstance(action, dict) and str(action.get("action") or "").strip()
            ],
            "risks": structured_risks,
        },
        "decision_packet": packet,
        "decision_packet_version": export_meta.get("schema_version", "1.0"),
    }
    return legacy


# --- PHASE 4: SYNTHESIS (The Extractor) ---
def synthesize_results(context, divergence_analysis=None, arbiter_report=None, results=None, user_id=None):
    """
    Extracts structured data (JSON) from the conversation history.
    Now includes divergence analysis for calibrated synthesis.
    Includes truth propagation: claim integrity is injected so synthesis
    avoids stale/unverified claims in downstream sections.
    """
    # Retrieve Workflow DNA for structure
    dna = WORKFLOW_DNA.get(context.workflow, WORKFLOW_DNA["RESEARCH"])
    if "alias" in dna:
        dna = WORKFLOW_DNA.get(dna["alias"], dna)

    # Flag assembly workflows — their final card renders the full document, not a summary
    SYNTHESIS_ASSEMBLY_WORKFLOWS = {"EOM_STATEMENT", "FINANCE", "AUDIT", "CODE_AUDIT", "LEGAL", "PORTFOLIO_BUILDER"}
    is_synthesis_assembly = context.workflow in SYNTHESIS_ASSEMBLY_WORKFLOWS

    # Build history with enriched context so the synthesizer knows exactly who did what
    history_text = ""
    for i, entry in enumerate(context.history):
        # We need to resolve the phase title for the synthesizer
        if context.workflow in WORKFLOW_PHASE_OVERRIDES:
            phase_set = WORKFLOW_PHASE_OVERRIDES[context.workflow]
            phase_title = phase_set.get(i, phase_set.get(min(i, 4))).get('title', f"Phase {i}")
        else:
            phase_title = f"Phase {i}"
        
        history_text += f"\n### PHASE {i}: {phase_title}\n"
        history_text += f"AGENT: {entry['ai'].upper()} as {entry['persona'].upper()}\n"
        # Cap each phase response to ~4000 chars to prevent synthesis prompt overflow
        phase_response = entry.get('response', '') or ''
        if len(phase_response) > 4000:
            phase_response = phase_response[:4000] + "\n[... TRUNCATED FOR SYNTHESIS ...]"
        history_text += f"RESPONSE:\n{phase_response}\n"

    # Inject divergence context if available
    if divergence_analysis and divergence_analysis.get('contested_topics'):
        history_text += "\n\n[ANALYTIC DIVERGENCE LAYER — PRE-SYNTHESIS CONTEXT]:\n"
        history_text += f"Consensus Score: {divergence_analysis.get('consensus_score', 'N/A')}/100\n"
        history_text += f"Divergence Score: {divergence_analysis.get('divergence_score', 'N/A')}/100\n"
        history_text += f"Summary: {divergence_analysis.get('divergence_summary', '')}\n"
        for topic in divergence_analysis.get('contested_topics', []):
            history_text += f"\nCONTESTED: {topic.get('topic', '')} (Severity: {topic.get('severity', 'unknown')})\n"
            for pos in topic.get('positions', []):
                history_text += f"  - {pos.get('provider', '').upper()}: {pos.get('position', '')}\n"

    # Inject arbiter corrections if available (Section 8a — mathematical consistency)
    if arbiter_report and arbiter_report.get('narrative_directive'):
        history_text += "\n\n[FINAL ARBITER — NUMERICAL CONSISTENCY LAYER]:\n"
        history_text += f"Status: {arbiter_report.get('status', 'N/A')}\n"
        history_text += arbiter_report['narrative_directive'] + "\n"
        history_text += (
            "\nCRITICAL: The above numerical values have been mathematically validated. "
            "You MUST use the resolved values in your synthesis. Do NOT use conflicting "
            "values from individual providers if they contradict the arbiter's resolved figures.\n"
        )

    # --- TRUTH PROPAGATION: Inject claim integrity into synthesis context ---
    # Items 1, 5, 6, 7: confidence labels, lifecycle tracking, output alignment
    if results:
        confirmed_claims = []
        unverified_claims = []
        invalidated_claims = []
        for provider, res in results.items():
            if not isinstance(res, dict) or not res.get('verified_claims'):
                continue
            for claim in res['verified_claims']:
                conf_label = claim.get('confidence', 'MEDIUM')
                lifecycle = claim.get('lifecycle', 'original')
                entry = (
                    f"{provider.upper()}: \"{claim['claim'][:120]}\" "
                    f"[{conf_label} confidence] "
                    f"(contribution: {claim.get('contribution', 0):+d}, lifecycle: {lifecycle})"
                )
                if claim['status'] == 'CONFIRMED':
                    confirmed_claims.append(entry)
                elif claim['status'] == 'SUPPORTED':
                    confirmed_claims.append(entry)
                else:
                    unverified_claims.append(entry)
                if claim.get('violations'):
                    for v in claim['violations']:
                        invalidated_claims.append(f"{provider.upper()}: \"{claim['claim'][:80]}\" — {v[:120]}")

        if confirmed_claims or unverified_claims or invalidated_claims:
            history_text += "\n\n[TRUTH PROPAGATION — CLAIM INTEGRITY LAYER]:\n"
            history_text += "The following claims have been cross-verified across council members.\n"
            history_text += "You MUST respect these findings in ALL synthesis sections.\n\n"
            if confirmed_claims:
                history_text += "VERIFIED CLAIMS (safe to include — attach confidence label in output):\n"
                for c in confirmed_claims[:15]:
                    history_text += f"  ✓ {c}\n"
            if unverified_claims:
                history_text += "\nUNVERIFIED CLAIMS (hedge language required — do NOT state as fact):\n"
                for c in unverified_claims[:15]:
                    history_text += f"  ? {c}\n"
                history_text += "\nRULE: Any unverified claim that appears in the executive summary, "
                history_text += "risk matrix, or recommendations MUST use hedging language "
                history_text += "(e.g. 'preliminary data suggests', 'subject to confirmation', "
                history_text += "'estimated range'). Do NOT present unverified claims as established fact.\n"
            if invalidated_claims:
                history_text += "\nINVALIDATED / CHALLENGED CLAIMS (REMOVE or REFRAME — never state as-is):\n"
                for c in invalidated_claims[:10]:
                    history_text += f"  ✗ {c}\n"
                history_text += "\nRULE: These claims have been flagged with violations "
                history_text += "(unsupported probability, vague temporal language, PQC non-compliance, etc.). "
                history_text += "Do NOT repeat them as-is. Either reframe with proper caveats or omit entirely. "
                history_text += "If a probability was flagged, either remove the specific number or add the basis.\n"

            # OUTPUT ALIGNMENT ENFORCEMENT (Item 6)
            history_text += "\n[OUTPUT ALIGNMENT — MANDATORY]:\n"
            history_text += "Your final narrative MUST reflect the claim validation status above.\n"
            history_text += "For each major claim in your output, include its confidence level:\n"
            history_text += "  - HIGH confidence claims: state directly\n"
            history_text += "  - MEDIUM confidence claims: use 'indicates' or 'suggests'\n"
            history_text += "  - LOW confidence claims: use 'preliminary' or 'unconfirmed'\n"
            history_text += "The executive summary and recommendations sections are the FINAL WORD.\n"
            history_text += "They must align with the truth score and verification outcomes — no contradictions.\n"

    # Section-specific synthesis directives — tell the LLM exactly what each section should contain
    SECTION_DIRECTIVES = {
        "executive_summary": (
            "FORMAT: Situation → Decision → Rationale → Confidence. Write as flowing prose, NOT labeled fields. "
            "Open with 1-2 sentences on the situation and what trends indicate. "
            "State the decision as a clear sentence: 'Decision: [action].' "
            "Follow with 1-2 sentences of rationale explaining why early intervention matters. "
            "Close with a confidence statement using bands like 'Moderate-High' with brief justification. "
            "MAX: 5-6 sentences total. VOICE: Senior strategy consultant writing a decision memo. "
            "NO fabricated percentages. NO specific dollar amounts unless in user data. "
            "Use directional language: 'elevated latency', 'rising outage frequency', 'growing risk'."
        ),
        "key_signals": (
            "List 4-6 concrete indicators that support the decision. "
            "Format as a bullet list. Each signal must reference a real trend or condition from the discussion. "
            "NO invented numbers — use directional language: 'sustained increase', 'upward trend', 'growing gap'. "
            "NO generic filler like 'market conditions remain uncertain'. "
            "These must come from the council discussion, not fabricated."
        ),
        "system_context": (
            "Identify the core entities involved and their roles in 3-5 bullets. "
            "Who or what is involved? What is their function in this decision? "
            "Keep it structural, not narrative. No elaboration — just entity + role. "
            "Example: 'Regional Fiber Network — degrading infrastructure layer under financial constraint'"
        ),
        "scenario_analysis": (
            "Present 2-3 scenarios: Most Plausible, Most Dangerous, and (if applicable) Strategic Opportunity. "
            "Each scenario needs: a name, 2-3 sentence description, directional progression timeline "
            "(e.g., 'near-term', 'mid-term' — NOT month-by-month or day-by-day), and an Implication line. "
            "NO fabricated precision — do NOT invent specific percentages like '30-35%' or timeframes like '12-48 hours'. "
            "Use directional language: 'elevated latency levels', 'extended outages', 'significantly higher costs'. "
            "NO stacked extreme outcomes in one scenario."
        ),
        "critical_challenges": (
            "Stress-test the decision. Present in two groups: "
            "HIGH IMPACT (2-3 items) and SECONDARY CONSIDERATIONS (2-3 items). "
            "For each challenge: state what it is, why it matters, and what would resolve it. "
            "End with a one-sentence synthesis: do the uncertainties outweigh the risk of inaction?"
        ),
        "tradeoffs": (
            "Build a structured comparison using [STRUCTURED_TABLE] format. "
            "Columns: Dimension | Short-Term Option | Long-Term Option. "
            "Dimensions: Operational, Financial, Customer/User, Strategic. "
            "Each cell: concise statement (1-2 sentences max). No paragraphs in cells."
        ),
        "decision": (
            "State the decision in ONE clear sentence. Then provide rationale as 3-5 bullet points. "
            "NO hedging. NO 'it depends'. This is the final call. "
            "Format: 'Decision: [Clear action statement]' followed by rationale bullets. "
            "NO fabricated cost figures or ROI percentages."
        ),
        "action_priorities": (
            "Break execution into 3 HIGH-LEVEL phases:\n"
            "Immediate: 1-2 actions\n"
            "Near-Term: 1-2 actions\n"
            "Mid-Term: 1-2 actions\n"
            "Actions must be executive-level decisions, NOT detailed playbooks. "
            "NO tool/vendor lists. NO day-by-day or month-by-month timelines. NO staffing plans. "
            "NO specific date ranges like '0-30 days' — use 'Immediate', 'Near-Term', 'Mid-Term'. "
            "Example good actions: "
            "'Identify and stabilize highest-risk network segments', "
            "'Upgrade core infrastructure in high-traffic areas', "
            "'Improve monitoring and network resilience'"
        ),
        "confidence": (
            "State the confidence band: High, Moderate-High, Moderate, or Low. "
            "KEY ASSUMPTIONS: 3-5 bullet points on what must hold true. "
            "LIMITATIONS: 2-3 bullet points on data gaps or unresolved questions. "
            "Confidence must reflect actual data quality, not optimism. "
            "If financial details or baselines are missing, confidence cannot be 'High' — use 'Moderate-High' at most. "
            "If UNKNOWNs outnumber VERIFIED facts, confidence cannot exceed Medium and score cannot exceed 70."
        ),
    }
    default_directive = "2-3 concise paragraphs synthesizing the council's findings for this section. Use specific data from the discussion. No filler."
    schema_sections = {
        section.lower().replace(" ", "_"): SECTION_DIRECTIVES.get(section.lower().replace(" ", "_"), default_directive)
        for section in dna["output_structure"]
    }

    prompt = f"""
    You are an Intelligence Synthesis Engine. Your goal is to convert a raw AI council discussion into a strict JSON Decision Packet.
    
    ## MISSION CONTEXT:
    - Type: {context.workflow}
    - Posture: {dna['posture']}
    - Outcome Expected: {dna['goal']}

    ## CRITICAL RULES:
    1. EXCLUSIVITY: Only use the provided COUNCIL DISCUSSION.
    2. STRICT SCHEMA: You MUST output exactly the JSON structure provided below. Do NOT add extra keys like 'meta' or 'sections'.
    3. NO CONVERSATIONAL FLUFF: Output only the requested data.
    4. BALANCED DECISION LOGIC: Never recommend halting all operations without actionable alternatives.
    5. CLEAN REPORT MODE: Write in ONE voice — senior strategy consultant delivering a memo.
       - NEVER expose internal labels: no node numbers, no model names, no provider names.
       - NEVER use "the council found" — use "analysis shows".
       - STRIP all internal tags: [VERIFIED], [CRITICAL] from text strings.
       - NO FAKE PRECISION: Do not invent percentages or timelines not found in user data.
    6. DECISION DISCIPLINE: No role may hide uncertainty. If evidence is insufficient, state 'INSUFFICIENT EVIDENCE'.
    7. CONFIDENCE CAP: If UNKNOWNs > VERIFIED facts, confidence MUST be MEDIUM or LOW and score MUST be <= 70.
    8. MISSING-DATA CAP: If critical data is missing, confidence MUST be LOW or CONDITIONAL.
    9. NO UNSOURCED PROJECTIONS: No timeline, probability, or numeric projection may appear unless it is explicitly supported in verified_claims or evidence_trace. If support is absent, classify it as an unknown instead.
    10. DIAGNOSTIC-FIRST RULE: If the prompt lacks quantified baselines, trend data, or verified root-cause evidence, do NOT force a specific root-cause fix. Emit a diagnostic-first recommendation and keep the decision at CONDITIONAL GO or NO-GO until evidence improves.

    FIELD INSTRUCTIONS:
    - decision_headline: One clear actionable sentence. NO hedging.
    - executive_summary: Situation -> Decision -> Rationale -> Confidence. 4-6 sentences max. No labeled fields in the text.
    - go_no_go_call.decision: Must strictly be GO, NO-GO, or CONDITIONAL GO. If evidence is sparse or baselines are missing, use CONDITIONAL GO or NO-GO and direct leadership to diagnostics first.
    - confidence: Must explicitly label band (HIGH | MEDIUM | LOW), a 0-100 score, and a 1-sentence basis. If UNKNOWNs > VERIFIED facts, band cannot exceed MEDIUM and score cannot exceed 70.
    - risk_vectors: Present the critical challenges and failure modes. High impact challenges only.
    - assumptions: List 3-5 baseline assumptions.
    - unknowns: List 2-3 critical unknowns.
    - immediate_actions: Break execution down. What happens Monday morning?
    - alternatives_rejected: What options were considered and explicitly killed by the Integrator?
    - verified_claims: claims verified from tags.
    - evidence_trace: critical data points mapping to source details. Any timeline, probability, or numeric projection used anywhere else in the packet must be supported here or in verified_claims.

    COUNCIL DISCUSSION:
    {history_text}

    Return ONLY a single valid JSON object mapped exactly to this schema:
    {{
      "decision_headline": "string",
      "executive_summary": "string",
      "go_no_go_call": {{
        "decision": "GO | NO-GO | CONDITIONAL GO",
        "rationale": "string"
      }},
      "confidence": {{
        "band": "HIGH | MEDIUM | LOW",
        "score": 0,
        "basis": "string"
      }},
      "verified_claims": [
        {{ "claim": "string", "status": "VERIFIED | PARTIAL | FALSE | UNVERIFIED", "source_ref": "string" }}
      ],
      "risk_vectors": [
        {{ "title": "string", "severity": "CRITICAL | HIGH | MEDIUM | LOW", "impact_type": "string", "description": "string", "mitigation": "string" }}
      ],
      "assumptions": ["string"],
      "unknowns": ["string"],
      "immediate_actions": [
        {{ "action": "string", "owner": "string", "timeline": "string" }}
      ],
      "alternatives_rejected": [
        {{ "option": "string", "reason_rejected": "string" }}
      ],
      "evidence_trace": [
        {{ "point": "string", "support": "string" }}
      ],
      "export_metadata": {{
        "report_type": "{context.workflow}",
        "workflow": "{context.workflow}",
        "generated_at": "{datetime.now().isoformat()}",
        "schema_version": "1.0",
        "composite_truth_score": 85
      }}
    }}
    CONFIDENCE CALIBRATION — MANDATORY SCORE-BAND ENFORCEMENT:
    {_build_confidence_directive(results)}
    """
    
    try:
        # Extract models from context
        models_used = [f"{entry['ai']} ({entry['persona']})" for entry in context.history]
        print(f"[SYNTHESIS] Prompt size: {len(prompt)} chars (~{len(prompt)//4} tokens). Calling GPT-4o...")
        resp = call_openai_gpt4(prompt, "Synthesizer", user_id=user_id)
        if not resp.get('success'):
            print(f"[SYNTHESIS] OpenAI failed: {resp.get('response', 'Unknown error')}")
            return {"executive_summary": "Synthesis unavailable", "meta": {"models_used": [], "truth_score": 0}}
        content = resp['response'].replace('```json', '').replace('```', '').strip()
        print(f"[SYNTHESIS] GPT-4o responded ({len(content)} chars). Parsing JSON...")
        try:
            data = json.loads(content)
        except json.JSONDecodeError as je:
            print(f"[SYNTHESIS] JSON parse failed: {je}. Attempting repair...")
            # Try to find the last valid closing brace
            last_brace = content.rfind('}')
            if last_brace > 0:
                try:
                    data = json.loads(content[:last_brace + 1])
                    print("[SYNTHESIS] JSON repaired via truncation.")
                except json.JSONDecodeError:
                    print("[SYNTHESIS] Repair failed. Returning fallback.")
                    return {"meta": {"models_used": models_used, "summary": "Synthesis JSON was malformed", "workflow": context.workflow}, "sections": {}, "structured_data": {}}
            else:
                return {"meta": {"models_used": models_used, "summary": "Synthesis returned no valid JSON", "workflow": context.workflow}, "sections": {}, "structured_data": {}}

        if _is_decision_packet_shape(data):
            print("[SYNTHESIS] Decision Packet detected — adapting to legacy export shape")
            data = adapt_decision_packet_to_legacy_shape(data, workflow=context.workflow)

        # Inject metadata if missing or simplified
        if "meta" not in data:
            data["meta"] = {}
        data["meta"]["models_used"] = models_used

        # Deterministic council_contributors — built from pipeline data, not LLM output
        contributors = []
        total_phases = len(context.history)
        
        # Resolve the active phase directive set for this workflow
        active_phase_directives = WORKFLOW_PHASE_OVERRIDES.get(context.workflow, PHASE_DIRECTIVES)
        
        for idx, entry in enumerate(context.history):
            if idx == total_phases - 1 and total_phases >= 3:
                phase_title = active_phase_directives.get(4, {}).get('title', f"Phase {idx + 1}")
            elif total_phases <= 5:
                phase_title = active_phase_directives.get(idx, {}).get('title', f"Phase {idx + 1}")
            else:
                phase_idx = min(int(idx / total_phases * 5), 4)
                phase_title = active_phase_directives.get(phase_idx, {}).get('title', f"Phase {idx + 1}")

            # Extract a real contribution summary from the response text.
            # Strip tags, markdown, and leading/trailing whitespace.
            # Use the LLM's own first substantive sentence as the summary.
            raw_resp = entry.get("response") or ""
            import re as _re
            _clean = _re.sub(r'\[/?[A-Z_]+\]', '', raw_resp)          # strip [TAGS]
            _clean = _re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', _clean) # strip **bold**
            _clean = _re.sub(r'^#+\s+', '', _clean, flags=_re.MULTILINE) # strip ## headers
            _clean = _clean.strip()
            # Find first sentence that is meaningful (>40 chars)
            _sentences = [s.strip() for s in _re.split(r'(?<=[.!?])\s+', _clean) if len(s.strip()) > 40]
            if _sentences:
                _summary = _sentences[0][:220]
                if len(_sentences[0]) > 220:
                    _summary += "..."
            else:
                _summary = _clean[:220].strip() or f"{phase_title} analysis completed."

            contributors.append({
                "phase": phase_title,
                "provider": entry["ai"],
                "role": entry.get("persona", ""),
                "contribution_summary": _summary,
            })
        data["council_contributors"] = contributors
        print(f"[SYNTHESIS] Injected {len(contributors)} council_contributors: {[c['phase'] for c in contributors]}")

        # DETERMINISTIC final_document injection for assembly workflows.
        # The LLM synthesizer is unreliable for this — pull the CFO/final phase
        # response directly from pipeline history instead of trusting the JSON field.
        SYNTHESIS_ASSEMBLY_WORKFLOWS = {"EOM_STATEMENT", "FINANCE", "AUDIT", "CODE_AUDIT", "LEGAL", "PORTFOLIO_BUILDER"}
        if context.workflow in SYNTHESIS_ASSEMBLY_WORKFLOWS and context.history:
            final_phase_response = context.history[-1].get("response", "").strip()
            if final_phase_response:
                # Strip raw intelligence tags so the display layer renders cleanly
                import re as _re2
                clean_doc = _re2.sub(r'\[DECISION_CANDIDATE\]([\s\S]*?)\[/DECISION_CANDIDATE\]', r'\1', final_phase_response)
                clean_doc = _re2.sub(r'\[RISK_VECTOR\]([\s\S]*?)\[/RISK_VECTOR\]', r'\1', clean_doc)
                clean_doc = _re2.sub(r'\[METRIC_ANCHOR\]([\s\S]*?)\[/METRIC_ANCHOR\]', r'\1', clean_doc)
                clean_doc = _re2.sub(r'\[TRUTH_BOMB\]([\s\S]*?)\[/TRUTH_BOMB\]', r'\1', clean_doc)
                clean_doc = _re2.sub(r'\[/?(?:DECISION_CANDIDATE|RISK_VECTOR|METRIC_ANCHOR|TRUTH_BOMB)\]', '', clean_doc)
                if "meta" not in data:
                    data["meta"] = {}
                data["meta"]["final_document"] = clean_doc
                print(f"[SYNTHESIS] Injected final_document from phase history ({len(clean_doc)} chars)")

        # NEW: Post-processing safety parse for tags (in case LLM misses some)
        import re
        decisions = re.findall(r'\[DECISION_CANDIDATE\](.*?)\[/DECISION_CANDIDATE\]', history_text, re.DOTALL)
        risks = re.findall(r'\[RISK_VECTOR\](.*?)\[/RISK_VECTOR\]', history_text, re.DOTALL)
        metrics = re.findall(r'\[METRIC_ANCHOR\](.*?)\[/METRIC_ANCHOR\]', history_text, re.DOTALL)
        
        # Merge safety-extracted tags if they aren't already in the JSON
        if not data.get("intelligence_tags"): data["intelligence_tags"] = {"decisions": [], "risks": [], "metrics": []}
        
        for d in decisions: 
            if d.strip() not in data["intelligence_tags"]["decisions"]: 
                data["intelligence_tags"]["decisions"].append(d.strip())
        for r in risks:
            if r.strip() not in data["intelligence_tags"]["risks"]:
                data["intelligence_tags"]["risks"].append(r.strip())
        for m in metrics:
            if m.strip() not in data["intelligence_tags"]["metrics"]:
                data["intelligence_tags"]["metrics"].append(m.strip())

        # --- POST-SYNTHESIS OUTPUT FILTER (Clean Report Mode) ---
        # Regex scrub on all sections to strip agent labels, tags, citations, tool spam
        sections = data.get("sections") or {}
        if sections:
            import re as _cf
            _clean_patterns = [
                (_cf.compile(r'NODE\s+\d+\s*[-—:]?\s*', _cf.IGNORECASE), ''),                  # NODE 01 —
                (_cf.compile(r'\b(OPENAI|ANTHROPIC|GOOGLE|PERPLEXITY|MISTRAL)\b', _cf.IGNORECASE), ''),  # provider names
                (_cf.compile(r'\b(GPT-4o?|Claude\s*\d*\.?\d*\s*\w*|Gemini\s*\d*\.?\d*|Sonar)\b', _cf.IGNORECASE), ''),  # model names
                (_cf.compile(r'\b(ANALYST|ARCHITECT|CRITIC|INTEGRATOR|COMPOSER)\b'), ''),        # phase labels
                (_cf.compile(r'\[(?:VERIFIED|CRITICAL|ACTION REQUIRED|REJECTED|AMBER|WARNING|RISK)\]'), ''),  # signal tags
                (_cf.compile(r'\[/?(?:TRUTH_BOMB|RISK_VECTOR|DECISION_CANDIDATE|METRIC_ANCHOR)\]'), ''),  # internal/closing tags
                (_cf.compile(r'\[\d+\]'), ''),                                                    # citation refs [1][2]
                (_cf.compile(r'KO-INT-\d+', _cf.IGNORECASE), ''),                                # session IDs
                (_cf.compile(r'(?:the\s+)?council\s+(?:found|determined|agreed|concluded)', _cf.IGNORECASE), 'analysis indicates'),  # council refs
                (_cf.compile(r'this\s+model\s+(?:said|found|concluded)', _cf.IGNORECASE), 'analysis shows'),  # model refs
            ]
            _scrubbed = 0
            for sec_key, sec_text in sections.items():
                if not isinstance(sec_text, str):
                    continue
                original = sec_text
                for pattern, replacement in _clean_patterns:
                    sec_text = pattern.sub(replacement, sec_text)
                # Clean up double spaces from removals
                sec_text = _cf.sub(r'  +', ' ', sec_text)
                sec_text = _cf.sub(r'\n +\n', '\n\n', sec_text)
                if sec_text != original:
                    _scrubbed += 1
                sections[sec_key] = sec_text
            data["sections"] = sections
            if _scrubbed:
                print(f"[OUTPUT FILTER] Scrubbed {_scrubbed} section(s) — removed agent labels, tags, citations")

        # Also scrub meta.summary if present
        meta_summary = (data.get("meta") or {}).get("summary", "")
        if meta_summary and isinstance(meta_summary, str):
            for pattern, replacement in _clean_patterns:
                meta_summary = pattern.sub(replacement, meta_summary)
            meta_summary = _cf.sub(r'  +', ' ', meta_summary)
            data["meta"]["summary"] = meta_summary

        # Clean report validation — log warnings (non-blocking)
        validate_clean_report(data, workflow=context.workflow)

        return data
    except Exception as e:
        print(f"[SYNTHESIS ERROR] {e}")
        return {"error": "Failed to synthesize structured data."}

# --- PHASE 6: ARTIFACT GENERATION (The Creator) ---
def generate_presentation_preview(synthesized_data, classification, user_id=None):
    """
    Generates a structured JSON outline for a presentation based on synthesized data.
    """
    output_type = classification.get('outputType', 'presentation')
    
    prompt = f"""
    You are a presentation designer. Create a slide deck outline from this data.

    SYNTHESIZED DATA:
    {json.dumps(synthesized_data, indent=2)}

    PRESENTATION TYPE: {output_type}
    DOMAIN: {classification.get('domain', 'business')}

    DESIGN PRINCIPLES (Mistral Mode - Section 10):
    1. Title slide + 1-2 slides per major section
    2. Max 5 bullets per slide (3-4 is better)
    3. VISUAL CREATIVITY (#1): Incorporate evocative visual descriptions (architecture diagrams, storyboards, 
       network maps). Focus on high-detail, non-generic aesthetics (avoid blue-teal-white cyberpunk overexposure).
    4. Executive-friendly language
    5. Clear narrative flow

    Return ONLY valid JSON (no markdown):
    {{
      "title": "Main Presentation Title",
      "subtitle": "Subtitle or context",
      "slideCount": 12,
      "estimatedDuration": "10-12 minutes",
      "slides": [
        {{
          "slideNumber": 1,
          "title": "Executive Summary",
          "content": [
            "Bullet 1",
            "Bullet 2"
          ],
          "visualType": "text|chart|diagram|image", 
          "chartData": {{ "type": "bar", "data": [] }},
          "layout": "title|content|two-column|image-text",
          "speakerNotes": "Notes for the presenter"
        }}
      ],
      "metadata": {{
        "theme": "professional|modern|minimal",
        "colorScheme": "blue-green"
      }}
    }}

    Generate { '12-15' if output_type == 'pitch_deck' else '8-12' } slides total.
    """
    
    try:
        # --- LEVERAGING MISTRAL STRENGTHS (#1 Visual Creativity & Multi-Tool Integration) ---
        # Using Mistral as the 'Designer' to generate the presentation outline.
        resp = call_mistral_api(prompt, "Designer", user_id=user_id)
        if not resp.get('success'):
             # Fallback to OpenAI if Mistral is unavailable
             resp = call_openai_gpt4(prompt, "Designer", user_id=user_id)
             
        content = resp['response'].replace('```json', '').replace('```', '').strip()
        return json.loads(content)
    except Exception as e:
        print(f"[PRESENTATION PREVIEW ERROR] {e}")
        return {"error": "Failed to generate presentation preview.", "slides": []}

# --- PRE-EXPORT VALIDATION (Clean Report Mode) ---

def validate_clean_report(synthesis_result, workflow="RESEARCH"):
    """
    Scan synthesis output for agent labels, verification tags, and tool spam.
    Non-blocking — returns warnings list for logging. Does not gate export.
    """
    warnings = []
    sections = synthesis_result.get("sections") or {}
    full_text = " ".join(str(v) for v in sections.values())

    # Agent label detection
    agent_patterns = [
        (r'NODE\s+\d+', "Node label"),
        (r'\b(OPENAI|ANTHROPIC|GOOGLE|PERPLEXITY|MISTRAL)\b', "Provider name"),
        (r'\b(GPT-4|Claude|Gemini|Sonar)\b', "Model name"),
    ]
    for pattern, label in agent_patterns:
        matches = re.findall(pattern, full_text, re.IGNORECASE)
        if matches:
            warnings.append(f"[CLEAN MODE] {label} found in synthesis: {matches[:3]}")

    # Verification tag detection
    tag_patterns = [
        (r'\[VERIFIED\]|\[CRITICAL\]|\[ACTION REQUIRED\]|\[REJECTED\]', "Signal tag"),
        (r'\[TRUTH_BOMB\]|\[RISK_VECTOR\]|\[DECISION_CANDIDATE\]|\[METRIC_ANCHOR\]', "Internal tag"),
    ]
    for pattern, label in tag_patterns:
        if re.search(pattern, full_text):
            warnings.append(f"[CLEAN MODE] {label} found in synthesis output")

    # Tool spam indicators
    if re.search(r'Day\s+\d+[\s:—-]', full_text):
        warnings.append("[CLEAN MODE] Day-by-day timeline detected in synthesis")
    if re.search(r'\[\d+\]', full_text):
        warnings.append("[CLEAN MODE] Citation references [N] detected in synthesis")

    if warnings:
        for w in warnings:
            print(w)
    else:
        print("[CLEAN MODE] Validation passed — no agent labels or tags detected")

    return warnings


# --- PHASE 7: ARTIFACT GENERATION (The Builder) ---
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def _sanitize_pptx_text(text):
    """Replace Unicode symbols that render as squares in default PowerPoint fonts."""
    import re
    if not isinstance(text, str):
        return str(text) if text else ""
    replacements = {
        '\u2022': '-', '\u2023': '-', '\u25aa': '-', '\u25cb': 'o',
        '\u2713': '[X]', '\u2714': '[X]', '\u2717': '[ ]', '\u2718': '[ ]',
        '\u2192': '->', '\u2190': '<-', '\u21d2': '=>',
        '\u2014': '--', '\u2013': '-',
        '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"',
        '\u2026': '...', '\u221e': 'inf', '\u2248': '~=', '\u00b1': '+/-',
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    text = re.sub(r'[\U0001F300-\U0001F9FF]', '', text)
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    return text


def generate_pptx_file(preview_data):
    """
    Builds a .pptx file from the preview JSON.
    Returns the file path.
    """
    prs = PPTXPresentation()

    # 1. Slide Master Layouts
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    two_col_layout = prs.slide_layouts[3] # Usually 'Two Content' or similar

    slides = preview_data.get('slides', [])

    for slide_data in slides:
        layout_type = slide_data.get('layout', 'content')

        # Select Layout
        if layout_type == 'title':
            slide = prs.slides.add_slide(title_layout)
        elif layout_type == 'two-column':
             slide = prs.slides.add_slide(two_col_layout)
        else:
            slide = prs.slides.add_slide(bullet_layout)

        # Set Title
        title = slide.shapes.title
        if title:
            title.text = _sanitize_pptx_text(slide_data.get('title', 'Untitled Slide'))

        # Set Content (Bullets)
        content_items = [_sanitize_pptx_text(item) for item in slide_data.get('content', [])]

        # Handle Layout Specifics
        if layout_type == 'title':
             # Subtitle is usually the second placeholder
             if len(slide.placeholders) > 1 and content_items:
                 subtitle = slide.placeholders[1]
                 subtitle.text = "\n".join(content_items)

        elif layout_type == 'two-column':
             # Left Column (Text)
             if len(slide.placeholders) > 1:
                 tf = slide.placeholders[1].text_frame
                 tf.text = content_items[0] if content_items else ""
                 for item in content_items[1:]:
                     p = tf.add_paragraph()
                     p.text = item
                     p.level = 0

             # Right Column (Chart Placeholder or Text)
             if len(slide.placeholders) > 2:
                 tf2 = slide.placeholders[2].text_frame
                 if slide_data.get('chartData'):
                     tf2.text = f"[CHART: {slide_data['chartData'].get('type', 'Chart')}]"
                 else:
                     tf2.text = ""

        else: # Standard Bullet Content
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = content_items[0] if content_items else ""
                for item in content_items[1:]:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0

        # Set font on all text frames to ensure consistent rendering
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(14)

        # Speaker Notes
        if slide_data.get('speakerNotes'):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = _sanitize_pptx_text(slide_data['speakerNotes'])

    # Save File
    filename = f"korum_artifact_{int(datetime.now().timestamp())}.pptx"
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)

    return filename # Return relative name for download


